import os
import math
import numpy as np
from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def generate_notebooklm_assets():
    width, height = 1920, 1080
    bg_dir = r"E:\Vaccine\ppt_assets_nlm"
    os.makedirs(bg_dir, exist_ok=True)

    # 1. 封面背景：溫潤紙質米白漸層 + 極細優雅網格 + 右側柔和光暈
    cover_img = Image.new("RGBA", (width, height), (250, 249, 246, 255))
    d_cover = ImageDraw.Draw(cover_img)
    
    # 雅緻極淡漸層
    for y in range(height):
        r = int(250 + (244 - 250) * (y / height))
        g = int(249 + (246 - 249) * (y / height))
        b = int(246 + (240 - 246) * (y / height))
        d_cover.line([(0, y), (width, y)], fill=(r, g, b, 255))
        
    # 極細淺灰色點陣網格 (NotebookLM 經典筆記網格)
    for gx in range(0, width, 48):
        for gy in range(0, height, 48):
            d_cover.ellipse([gx-1, gy-1, gx+1, gy+1], fill=(226, 232, 240, 140))

    # 右側柔和彩色光暈 (Google/NotebookLM 標誌性柔光)
    glow_overlay = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    d_glow = ImageDraw.Draw(glow_overlay)
    
    # 柔和紫藍/薄荷綠漸變光圈
    d_glow.ellipse([width - 650, -100, width + 250, 750], fill=(219, 234, 254, 120)) # soft blue
    d_glow.ellipse([width - 450, 200, width + 350, 950], fill=(204, 251, 241, 100)) # soft teal
    d_glow.ellipse([width - 750, 450, width - 50, 1150], fill=(254, 243, 199, 90))  # soft amber

    glow_blur = glow_overlay.filter(ImageFilter.GaussianBlur(radius=80))
    cover_final = Image.alpha_composite(cover_img, glow_blur)
    
    cover_path = os.path.join(bg_dir, "nlm_cover_bg.png")
    cover_final.convert("RGB").save(cover_path, quality=95)

    # 2. 內頁背景：純淨極簡柔白 + 頂部極輕微呼吸光
    inner_img = Image.new("RGBA", (width, height), (252, 251, 249, 255))
    d_inner = ImageDraw.Draw(inner_img)
    for y in range(height):
        r = int(253 + (248 - 253) * (y / height))
        g = int(252 + (249 - 252) * (y / height))
        b = int(250 + (246 - 250) * (y / height))
        d_inner.line([(0, y), (width, y)], fill=(r, g, b, 255))
        
    for gx in range(0, width, 48):
        for gy in range(0, height, 48):
            d_inner.ellipse([gx-1, gy-1, gx+1, gy+1], fill=(226, 232, 240, 100))

    inner_path = os.path.join(bg_dir, "nlm_inner_bg.png")
    inner_img.convert("RGB").save(inner_path, quality=95)

    return cover_path, inner_path

def build_notebooklm_presentation():
    cover_bg, inner_bg = generate_notebooklm_assets()
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # NotebookLM 頂級美學調色盤 (Warm Minimalist & Editorial)
    C_CARD_BG = RGBColor(255, 255, 255)         # #FFFFFF 純白質感卡片
    C_CARD_BG_ALT = RGBColor(248, 250, 252)     # #F8FAFC 微灰白對比卡片
    C_BORDER = RGBColor(226, 232, 240)          # #E2E8F0 極細淺灰邊框
    C_BORDER_ACTIVE = RGBColor(14, 165, 233)    # #0EA5E9 聚焦藍
    
    C_TITLE = RGBColor(15, 23, 42)              # #0F172A 深沉雅緻石墨黑
    C_BODY = RGBColor(51, 65, 85)               # #334155 溫潤高讀性炭灰
    C_MUTED = RGBColor(100, 116, 139)           # #64748B 質感輔助灰
    
    C_BLUE = RGBColor(2, 132, 199)              # #0284C7 品牌湛藍
    C_TEAL = RGBColor(13, 148, 136)             # #0D9488 薄荷深青
    C_CORAL = RGBColor(225, 29, 72)             # #E11D48 亮點珊瑚紅
    C_AMBER = RGBColor(217, 119, 6)             # #D97706 琥珀棕

    def add_base_slide(title_text, subtitle_text=None, category="NOTEBOOKLM STUDY GUIDE"):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(inner_bg, 0, 0, prs.slide_width, prs.slide_height)

        # 頂部 NotebookLM 膠囊標籤
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.42), Inches(3.2), Inches(0.32))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(241, 245, 249)
        badge.line.color.rgb = RGBColor(203, 213, 225)
        badge.line.width = Pt(0.75)

        tb_b = slide.shapes.add_textbox(Inches(0.8), Inches(0.42), Inches(3.2), Inches(0.32))
        tf_b = tb_b.text_frame
        tf_b.margin_left = Inches(0.12)
        tf_b.margin_top = Inches(0.04)
        p_b = tf_b.paragraphs[0]
        p_b.text = f"✦ {category}"
        p_b.font.name = "Segoe UI"
        p_b.font.size = Pt(9.5)
        p_b.font.bold = True
        p_b.font.color.rgb = C_BLUE

        # 標題文字方塊
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.82), Inches(11.7), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Microsoft JhengHei"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = C_TITLE

        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = "Microsoft JhengHei"
            p2.font.size = Pt(12)
            p2.font.color.rgb = C_MUTED
            p2.space_before = Pt(3)

        return slide

    def add_card(slide, left, top, width, height, title, content_list, badge=None, bg_color=C_CARD_BG, border_color=C_BORDER, header_color=C_BLUE):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.25)

        tb = slide.shapes.add_textbox(left + Inches(0.24), top + Inches(0.22), width - Inches(0.48), height - Inches(0.44))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Microsoft JhengHei"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = header_color

        if badge:
            p.text = f"{title}  ·  {badge}"

        for item in content_list:
            p_item = tf.add_paragraph()
            p_item.text = item
            p_item.font.name = "Microsoft JhengHei"
            p_item.font.size = Pt(11.5)
            p_item.font.color.rgb = C_BODY
            p_item.space_before = Pt(5)
            if item.startswith("•") or item.startswith("👉") or item.startswith("✓") or item.startswith("🚨") or item.startswith("【") or item.startswith("1.") or item.startswith("2.") or item.startswith("3."):
                p_item.font.bold = False
                p_item.font.color.rgb = C_TITLE

    # -------------------------------------------------------------
    # 投影片 1：封面 (Cover Slide - NotebookLM Style)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    slide1.shapes.add_picture(cover_bg, 0, 0, prs.slide_width, prs.slide_height)

    # 封面頂部微型標籤
    tag_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.4), Inches(3.4), Inches(0.36))
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = RGBColor(238, 242, 255)
    tag_box.line.color.rgb = RGBColor(199, 210, 254)
    tag_box.line.width = Pt(1)
    
    t_tag = tag_box.text_frame
    p_t = t_tag.paragraphs[0]
    p_t.text = "✦ AI-NATIVE SOFTWARE CREATION"
    p_t.font.name = "Segoe UI"
    p_t.font.size = Pt(10)
    p_t.font.bold = True
    p_t.font.color.rgb = C_BLUE

    tb1 = slide1.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.5), Inches(4.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "何謂 Vibe Coding？"
    p1.font.name = "Microsoft JhengHei"
    p1.font.size = Pt(46)
    p1.font.bold = True
    p1.font.color.rgb = C_TITLE

    p2 = tf1.add_paragraph()
    p2.text = "從「台灣預防接種指南助手」看 AI 時代人機共創的全新軟體開發典範"
    p2.font.name = "Microsoft JhengHei"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = C_BLUE
    p2.space_before = Pt(10)

    p3 = tf1.add_paragraph()
    p3.text = "實務專案案例剖析 ‧ 概念演進 ‧ 人機角色蛻變 ‧ 醫療輔助系統四端跨平台落地"
    p3.font.name = "Microsoft JhengHei"
    p3.font.size = Pt(13.5)
    p3.font.color.rgb = C_MUTED
    p3.space_before = Pt(20)

    p4 = tf1.add_paragraph()
    p4.text = "核心案例：Rust 醫療算法 ‧ WebAssembly 靜態網頁 ‧ Android 原生 APK ‧ iOS 原生日曆連動"
    p4.font.name = "Microsoft JhengHei"
    p4.font.size = Pt(11.5)
    p4.font.color.rgb = C_TEAL
    p4.space_before = Pt(8)

    # -------------------------------------------------------------
    # 投影片 2：什麼是 Vibe Coding？
    # -------------------------------------------------------------
    slide2 = add_base_slide(
        "1. 什麼是 Vibe Coding？核心定義與概念演進",
        "由前 Tesla AI 總監、OpenAI 共同創辦人 Andrej Karpathy 於 2025 年提出的顛覆性心智模型",
        category="CONCEPT DEFINITION"
    )
    add_card(
        slide2, Inches(0.8), Inches(1.9), Inches(5.6), Inches(5.0),
        "💡 Vibe Coding 的核心本質",
        [
            "• 開發者不再逐字手寫語法，而是以「直覺、意圖與自然語言對話」引導 AI。",
            "• 開發者的角色轉變為【總架構師】與【產品質檢官 (Director & Reviewer)】。",
            "• 重點放在「想要什麼體驗 (Vibe) 與核心業務邏輯」，讓 AI Agent 負責底層實作、修復報錯、編譯與跨平台部署。",
            "",
            "👉 『You just see stuff, say stuff, run stuff, and copy-paste stuff. It works.』"
        ],
        badge="全新心智模型",
        header_color=C_BLUE
    )
    add_card(
        slide2, Inches(6.8), Inches(1.9), Inches(5.7), Inches(5.0),
        "⚡ 傳統開發 vs. Vibe Coding 對比",
        [
            "【傳統工程模式 (Traditional Coding)】",
            "• 手動查閱 API 文件、記憶語言語法細節、手寫樣板代碼",
            "• 遭遇跨平台 (Web / Android / NDK) 編譯地獄時耗費數天排錯",
            "",
            "【Vibe Coding 人機協作模式】",
            "• 領域專家（如醫師/PM）直接以專業業務邏輯下達指令",
            "• AI 即刻進行法規文獻查證、Rust 算法撰寫、Wasm 封裝、APK 簽署",
            "• 敏捷迭代：一個複雜臨床需求在數分鐘內完成四端發布"
        ],
        badge="效率革命",
        border_color=RGBColor(167, 243, 208),
        header_color=C_TEAL
    )

    # -------------------------------------------------------------
    # 投影片 3：經典案例剖析 - 台灣預防接種指南助手
    # -------------------------------------------------------------
    slide3 = add_base_slide(
        "2. 實戰案例剖析：台灣預防接種指南助手",
        "以真實複雜的醫療系統為例——如何透過自然語言對話完成四端跨平台同步開發",
        category="CASE STUDY"
    )
    add_card(
        slide3, Inches(0.8), Inches(1.9), Inches(3.7), Inches(5.0),
        "🏥 臨床專業意圖",
        [
            "• 「診所需要算公自費疫苗時程」",
            "• 「7歲以上兒童有生長曲線嗎？請引用台灣最新官方年份」",
            "• 「要能支援 iOS 相機直掃行事曆」",
            "• 「家長改日期若提早，需跳出醫療法規理由警示」",
            "",
            "👉 專注於【臨床核心邏輯】，完全不需觸碰底層語法。"
        ],
        badge="Doctor's Intent",
        header_color=C_BLUE
    )
    add_card(
        slide3, Inches(4.8), Inches(1.9), Inches(3.7), Inches(5.0),
        "🤖 AI Agent 自主實現",
        [
            "• 聯網查核 2024 衛福部國健署標準",
            "• 撰寫嚴謹的 Rust 演算法 (vaccine-core)",
            "• 編譯 WebAssembly 驅動純前端運算",
            "• 自動配置 Android NDK 工具鏈打包簽署 APK",
            "• 生成 iCalendar (.ics) / VEVENT 數據流"
        ],
        badge="Execution",
        border_color=RGBColor(186, 230, 253),
        header_color=C_TEAL
    )
    add_card(
        slide3, Inches(8.8), Inches(1.9), Inches(3.7), Inches(5.0),
        "🚀 四端跨平台落地",
        [
            "1. 🌐 GitHub Pages 靜態無伺服器網站",
            "2. 📱 Android 原生 arm64 簽署 APK",
            "3. 💻 Windows 單一免安裝獨立綠色版",
            "4. 🍏 iOS 相機直掃 iCal 原生日曆",
            "",
            "✨ 從單一需求到四端發布，全由對話驅動！"
        ],
        badge="Deliverables",
        border_color=RGBColor(254, 215, 170),
        header_color=C_AMBER
    )

    # -------------------------------------------------------------
    # 投影片 4：Vibe Coding 運作流程
    # -------------------------------------------------------------
    slide4 = add_base_slide(
        "3. Vibe Coding 的四步驟高頻迭代循環",
        "從自然語言需求到跨平台生產環境交付的完整閉環 (Feedback Loop)",
        category="ITERATIVE WORKFLOW"
    )
    steps = [
        ("Step 1: 意圖表達", "使用者提出自然語言需求\n例如：「增加iOS行事曆QR Code與自訂日期」", C_BLUE),
        ("Step 2: 查核與架構", "AI 搜尋法規、規劃資料模型\n設計前後端跨語言序列化格式", C_TEAL),
        ("Step 3: 實現與編譯", "AI 撰寫程式碼、呼叫編譯器\n遇到報錯自主排查修復 (Self-Healing)", C_CORAL),
        ("Step 4: 驗證與發布", "產出 Android APK、推送到 GitHub\n使用者即刻在手機/網頁體驗", C_AMBER),
    ]
    for i, (title, desc, color) in enumerate(steps):
        left = Inches(0.8 + i * 2.95)
        add_card(
            slide4, left, Inches(2.0), Inches(2.8), Inches(4.8),
            title,
            [desc, "", "👉 立即進入下一輪反饋循環 (Feedback Loop)"],
            header_color=color
        )

    # -------------------------------------------------------------
    # 投影片 5：關鍵支柱與成功要素
    # -------------------------------------------------------------
    slide5 = add_base_slide(
        "4. Vibe Coding 成功落地的三大支柱",
        "為什麼本專案能做到又快又穩？底層的技術架構與工程支撐",
        category="KEY ENABLERS"
    )
    add_card(
        slide5, Inches(0.8), Inches(1.9), Inches(3.7), Inches(5.0),
        "1. Skill 與上下文沉澱",
        [
            "• 建立專屬 SKILL.md 開發規範",
            "• 將 NDK 環境變數、Wasm 序列化規範、Git 部署指令結構化",
            "• 讓 AI 每次對話都能維持統一水準，不遺忘專案背景知識"
        ],
        badge="Memory",
        header_color=C_BLUE
    )
    add_card(
        slide5, Inches(4.8), Inches(1.9), Inches(3.7), Inches(5.0),
        "2. 工具調用能力",
        [
            "• 具備終端機 (PowerShell)、網路搜尋、檔案精確替換工具",
            "• 遇到 Android 編譯錯誤時，AI 能自主讀取報錯並修正 Rust / Gradle 配置",
            "• 真正的【自主除錯與修復 (Self-Healing)】"
        ],
        badge="Tool Calling",
        border_color=RGBColor(186, 230, 253),
        header_color=C_TEAL
    )
    add_card(
        slide5, Inches(8.8), Inches(1.9), Inches(3.7), Inches(5.0),
        "3. 領域專業引導",
        [
            "• 人類的不可替代性在於【臨床專業判斷】",
            "• 審核疫苗間隔邏輯是否合規、評估使用者互動體驗",
            "• 領域專家 + 強大 AI = 專業級系統迅速誕生"
        ],
        badge="Human-in-the-Loop",
        border_color=RGBColor(254, 215, 170),
        header_color=C_AMBER
    )

    # -------------------------------------------------------------
    # 投影片 6：總結與未來展望
    # -------------------------------------------------------------
    slide6 = add_base_slide(
        "5. 總結：軟體開發的文藝復興時代",
        "每個人都能成為自己領域的數位架構師 (Domain Architects)",
        category="CONCLUSION"
    )
    add_card(
        slide6, Inches(0.8), Inches(1.9), Inches(5.6), Inches(5.0),
        "🚀 核心啟示 (Key Takeaways)",
        [
            "1. 程式碼不再是壁壘：思想與領域知識 (Domain Knowledge) 才是核心價值。",
            "2. 敏捷的極致體現：從醫療想法到多端上線僅需數小時，大幅降低創新試錯成本。",
            "3. 架構思維 > 語法記憶：掌握系統設計與需求拆解能力，比死記 API 更有競爭力。"
        ],
        badge="Takeaways",
        header_color=C_BLUE
    )
    add_card(
        slide6, Inches(6.8), Inches(1.9), Inches(5.7), Inches(5.0),
        "🌟 未來展望 (The Future of Creation)",
        [
            "• 醫生可以為診所打造最合身的醫療輔助系統",
            "• 老師可以為學生建構客製化的互動教學軟體",
            "• 創業家可以在一天內驗證並發布 MVP 產品",
            "",
            "✨ 『在 AI 時代，你的想像力與專業知識，就是最強大的程式語言。』"
        ],
        badge="Vision",
        border_color=RGBColor(254, 215, 170),
        header_color=C_AMBER
    )

    output_path = r"E:\Vaccine\Vibe_Coding_Presentation.pptx"
    prs.save(output_path)
    print(f"NotebookLM style presentation generated: {output_path}")

if __name__ == "__main__":
    build_notebooklm_presentation()
