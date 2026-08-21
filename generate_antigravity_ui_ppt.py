import os
from PIL import Image, ImageDraw, ImageFont, ImageFilter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def generate_notebooklm_assets():
    width, height = 1920, 1080
    bg_dir = r"E:\Vaccine\ppt_assets_nlm"
    os.makedirs(bg_dir, exist_ok=True)

    # 1. 封面背景：溫潤紙質米白漸層 + 極細優雅網格 + 右側柔和彩色光暈
    cover_img = Image.new("RGBA", (width, height), (250, 249, 246, 255))
    d_cover = ImageDraw.Draw(cover_img)
    
    for y in range(height):
        r = int(250 + (244 - 250) * (y / height))
        g = int(249 + (246 - 249) * (y / height))
        b = int(246 + (240 - 246) * (y / height))
        d_cover.line([(0, y), (width, y)], fill=(r, g, b, 255))
        
    for gx in range(0, width, 48):
        for gy in range(0, height, 48):
            d_cover.ellipse([gx-1, gy-1, gx+1, gy+1], fill=(226, 232, 240, 140))

    glow_overlay = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    d_glow = ImageDraw.Draw(glow_overlay)
    
    d_glow.ellipse([width - 650, -100, width + 250, 750], fill=(219, 234, 254, 120)) # soft blue
    d_glow.ellipse([width - 450, 200, width + 350, 950], fill=(204, 251, 241, 100)) # soft teal
    d_glow.ellipse([width - 750, 450, width - 50, 1150], fill=(254, 243, 199, 90))  # soft amber

    glow_blur = glow_overlay.filter(ImageFilter.GaussianBlur(radius=80))
    cover_final = Image.alpha_composite(cover_img, glow_blur)
    
    cover_path = os.path.join(bg_dir, "nlm_cover_bg.png")
    cover_final.convert("RGB").save(cover_path, quality=95)

    # 2. 內頁背景：純淨極簡柔白 + 網格底紋
    inner_img = Image.new("RGBA", (width, height), (252, 251, 249, 255))
    d_inner = ImageDraw.Draw(inner_img)
    for y in range(height):
        r = int(253 + (248 - 253) * (y / height))
        g = int(252 + (249 - 252) * (y / height))
        b = int(250 + (246 - 250) * (y / height))
        d_inner.line([(0, y), (width, y)], fill=(r, g, b, 255))
        
    for gx in range(0, width, 48):
        for gy in range(0, height, 48):
            d_inner.ellipse([gx-1, gy-1, gx+1, gy+1], fill=(226, 232, 240, 100))

    inner_path = os.path.join(bg_dir, "nlm_inner_bg.png")
    inner_img.convert("RGB").save(inner_path, quality=95)

    return cover_path, inner_path

# 1. 建立高解析度 Antigravity UI 精緻介面 Mockup 截圖 (800x480)
def create_antigravity_ui_mockups():
    img_dir = r"E:\Vaccine\ppt_assets_ui"
    os.makedirs(img_dir, exist_ok=True)
    
    # 配色常數 (Antigravity IDE Dark Palette)
    C_IDE_BG = (15, 23, 42)          # #0F172A 深沉科技深藍
    C_IDE_SIDEBAR = (30, 41, 59)     # #1E293B 側邊欄/卡片底色
    C_IDE_HEADER = (15, 23, 42)      # #0F172A
    C_BORDER = (51, 65, 85)          # #334155
    C_TEXT_WHITE = (248, 250, 252)
    C_TEXT_MUTED = (148, 163, 184)
    C_ACCENT_BLUE = (56, 189, 248)   # #38BDF8
    C_ACCENT_TEAL = (45, 212, 191)   # #2DD4BF
    C_ACCENT_AMBER = (251, 191, 36)  # #FBBF24
    C_ACCENT_CORAL = (248, 113, 113) # #F87171
    C_ACCENT_PURPLE = (192, 132, 252)# #C084FC

    # 嘗試載入中文字型
    font_path = "C:/Windows/Fonts/msjh.ttc"
    font_bold_path = "C:/Windows/Fonts/msjhbd.ttc"
    
    try:
        font_title = ImageFont.truetype(font_bold_path, 18)
        font_sub = ImageFont.truetype(font_path, 13)
        font_code = ImageFont.truetype(font_path, 12)
        font_badge = ImageFont.truetype(font_bold_path, 11)
    except:
        font_title = font_sub = font_code = font_badge = ImageFont.load_default()

    mockups = [
        {
            "id": "ui_p2_intent",
            "title": "Antigravity Chat · Natural Language Intent",
            "badge": "PROMPT AS LOGIC",
            "badge_col": C_ACCENT_BLUE,
            "user_prompt": "醫師提問：「病人不見得當時有空，可以修改日期嗎？如果提早要顯示醫療理由」",
            "action_title": "Antigravity Agent 正在規劃並執行任務：",
            "steps": [
                "1. 讀取 `vaccine-app/src/index.html` 內的 Calendar Modal 結構",
                "2. 調用 replace_file_content 加入 <input type='date'> 與重設按鈕",
                "3. 實作 validateAndExplainDate() 即時評估提早/延後醫療理由",
                "4. 即時同步 Google 日曆 URL 與 iOS VEVENT QR Code"
            ],
            "status": "STATUS: COMPLETED IN 8.2s · LIVE UPDATED"
        },
        {
            "id": "ui_p3_tools",
            "title": "Antigravity Engine · Tool Calling Pipeline",
            "badge": "AUTONOMOUS TOOLS",
            "badge_col": C_ACCENT_TEAL,
            "user_prompt": "Antigravity 自主工具調用鏈 (Tool Call Sequence)",
            "action_title": "Active Tools in Action:",
            "steps": [
                "🛠️ tool_call: replace_file_content(path='vaccine-app/src/main.js')",
                "🛠️ tool_call: run_command(cmd='cargo build --target aarch64-linux-android')",
                "🛠️ tool_call: search_web(query='Taiwan MOHW 0-18 BMI standard 2024')",
                "🛠️ tool_call: invoke_subagent(role='CSS Eye-Friendly Reskin')"
            ],
            "status": "PIPELINE: 4 TOOLS CALLED CONCURRENTLY · ZERO MANUAL CODE"
        },
        {
            "id": "ui_p4_medical",
            "title": "Antigravity Logic · Rust Clinical State Machine",
            "badge": "RUST VACCINE-CORE",
            "badge_col": C_ACCENT_TEAL,
            "user_prompt": "衛福部 CDC 接種時程與 ACIP 補打間隔演算法轉化",
            "action_title": "vaccine-core/src/lib.rs (State Machine):",
            "steps": [
                "• pub fn calculate_catch_up(vaccine_id, last_dose, last_date) -> String",
                "• Match rule: 五合一 D4 需滿 180 天 (6個月) 且年齡 >= 18 個月",
                "• Match rule: PCV13 於 7~11 個月補打時自動縮減為 3 劑 (2+1)",
                "• Co-administration: 自動比對活性減毒同天或隔 28 天指引"
            ],
            "status": "STATUS: 100% COMPILED TO WASM · ZERO SERVER RISK"
        },
        {
            "id": "ui_p5_research",
            "title": "Antigravity Research · Web Citation Grounding",
            "badge": "SEARCH_WEB TOOL",
            "badge_col": C_ACCENT_AMBER,
            "user_prompt": "醫師質疑：「超過7歲生長曲線是哪一年資料？有更新的嗎？」",
            "action_title": "Antigravity 聯網查核衛生福利部國民健康署公報：",
            "steps": [
                "🌐 search_web: '衛福部 國健署 兒童生長曲線 7-18歲 BMI 百分位'",
                "📄 查證結果：2009 銜接標準、2010 衛署授升字第0990700680號公告",
                "✓ 確認現行 2024《兒童健康手冊》仍以此為唯一法定標準",
                "🏷️ 在 index.html 與報告卡片頂部正式標註官方出處與年份"
            ],
            "status": "GROUNDING: 100% VERIFIED WITH MOHW OFFICIAL DATASETS"
        },
        {
            "id": "ui_p6_wasm",
            "title": "Antigravity Bridge · Wasm Zero-Server Architecture",
            "badge": "WASM BRIDGE",
            "badge_col": C_ACCENT_PURPLE,
            "user_prompt": "診所隱私保護：零伺服器傳輸純前端運算",
            "action_title": "Rust Wasm <-> JS Contract Architecture:",
            "steps": [
                "1. Rust 端：使用 serde_json::to_string 輸出標準 JSON 字串",
                "2. 前端 app.js：fallbackInvoke 接收後 JSON.parse(res) 還原",
                "3. 變數名稱映射：手動對齊 CamelCase <-> snake_case",
                "4. 離線執行：病人生日與健康數據 100% 留在本機瀏覽器"
            ],
            "status": "SECURITY: ZERO EXTERNAL SERVER CALLS · HIPAA COMPLIANT"
        },
        {
            "id": "ui_p7_ndk",
            "title": "Antigravity DevOps · Android NDK Self-Healing",
            "badge": "SELF-HEALING DEVOPS",
            "badge_col": C_ACCENT_CORAL,
            "user_prompt": "編譯 Android APK 遭遇 Clang Linker 255 與符號重複報錯",
            "action_title": "Antigravity 自主診斷並修復環境問題：",
            "steps": [
                "🚨 偵測錯誤：vaccine-core cdylib 與 rlib 產生交叉編譯符號衝突",
                "🔧 自動修復：修改 Cargo.toml 改為純 rlib，配置 Clang 34 旗標",
                "⚙️ run_command: cargo build --target aarch64-linux-android",
                "📦 自動調用: gradlew.bat assembleArm64Release & apksigner 簽署"
            ],
            "status": "OUTPUT: E:\\台灣疫苗指南助手.apk (SIGNED & READY)"
        },
        {
            "id": "ui_p8_ical",
            "title": "Antigravity UX · Dual Ecosystem Calendar Bridge",
            "badge": "DUAL ECOSYSTEM",
            "badge_col": C_ACCENT_BLUE,
            "user_prompt": "醫師需求：「生成行事曆 QR Code，有可能增加 iOS 日曆嗎？」",
            "action_title": "雙系統跨生態日曆協議實作：",
            "steps": [
                "🌐 Google 日曆：生成 calendar.google.com/render?action=TEMPLATE 網址",
                "🍏 iOS 檔案：生成標準 RFC-5545 iCalendar (.ics) Blob 一鍵下載",
                "📷 iOS 相機直掃：將 QR Code 切換為 BEGIN:VEVENT...END:VEVENT",
                "✓ iPhone 內建相機一照，頂部立即彈出「加入日曆」按鈕"
            ],
            "status": "COMPATIBILITY: IOS CAMERA DIRECT SCAN + GOOGLE CALENDAR"
        },
        {
            "id": "ui_p9_warning",
            "title": "Antigravity Clinical · Safety Interval Validation",
            "badge": "SAFETY INTERVALS",
            "badge_col": C_ACCENT_CORAL,
            "user_prompt": "自訂預約提醒日期 ➔ 動態觸發臨床提早/延後警示",
            "action_title": "validateAndExplainDate() 即時臨床安全評估：",
            "steps": [
                "🚨 提早接種 (紅色警示)：若早於最小年齡，因母體抗體干擾導致失效",
                "ℹ️ 延後接種 (黃色提示)：生病延後直接順延接種即可，不需從頭重打",
                "📅 即時聯動：修改日期時，Google URL 與 iOS QR Code 瞬間同步",
                "✓ 醫師思維邏輯無縫嵌入互動介面"
            ],
            "status": "CLINICAL LOGIC: DYNAMIC MEDICAL REASONING EMBEDDED"
        },
        {
            "id": "ui_p10_skills",
            "title": "Antigravity Memory · Project Skills Brain",
            "badge": "PERSISTENT SKILL",
            "badge_col": C_ACCENT_BLUE,
            "user_prompt": "知識沉澱：.agents/skills/vaccine-dev/SKILL.md",
            "action_title": "專案長期大腦記憶結構：",
            "steps": [
                "• 記錄 Android NDK / Clang 34 編譯環境變數與路徑",
                "• 記錄 Rust Wasm JSON 傳輸合約與 CamelCase 變數映射表",
                "• 記錄 0-18 歲生長曲線 2009/2010 國健署法定年份標準",
                "• 記錄 GitHub Pages 部署與快取強刷 SOP"
            ],
            "status": "MEMORY: PERSISTENT CONTEXT ACROSS LONG CONVERSATIONS"
        },
        {
            "id": "ui_p11_subagents",
            "title": "Antigravity Parallel · Specialized Subagents",
            "badge": "SUBAGENT WORKER",
            "badge_col": C_ACCENT_PURPLE,
            "user_prompt": "全系統色彩重構：切換為護眼溫潤大地色系 (Warm Linen)",
            "action_title": "invoke_subagent(role='CSS Eye-Friendly Reskin'):",
            "steps": [
                "👥 主代理保持高階對話，派發專屬 Subagent 平行重構 styles.css",
                "🎨 更新 20+ 組 CSS 變數：溫潤亞麻背景 + 低飽和青藍頁籤",
                "🏷️ 重新設計常規綠、補助藍、自費棕、高風險紅等標籤色",
                "✓ 子代理完成後自動回報編譯成功，不污染主對話 Context"
            ],
            "status": "PARALLEL: BACKGROUND REFACTORING COMPLETED SEAMLESSLY"
        }
    ]

    generated_paths = {}
    for m in mockups:
        img = Image.new("RGBA", (820, 480), C_IDE_BG)
        draw = ImageDraw.Draw(img)
        
        # 繪製外框與頂部視窗列
        draw.rectangle([(0, 0), (819, 479)], fill=C_IDE_BG, outline=C_BORDER, width=2)
        draw.rectangle([(0, 0), (819, 42)], fill=C_IDE_SIDEBAR)
        draw.line([(0, 42), (819, 42)], fill=C_BORDER, width=1)
        
        # 視窗紅黃綠控制按鈕
        draw.ellipse([(14, 15), (24, 25)], fill=(239, 68, 68))
        draw.ellipse([(32, 15), (42, 25)], fill=(245, 158, 11))
        draw.ellipse([(50, 15), (60, 25)], fill=(16, 185, 129))
        
        # 頂部標題
        draw.text((80, 12), m["title"], fill=C_TEXT_WHITE, font=font_title)
        
        # 右上角膠囊標籤
        badge_w = len(m["badge"]) * 8 + 20
        bx2 = 805
        bx1 = bx2 - badge_w
        draw.rounded_rectangle([(bx1, 10), (bx2, 32)], radius=4, fill=(20, 30, 48), outline=m["badge_col"], width=1)
        draw.text((bx1 + 10, 14), m["badge"], fill=m["badge_col"], font=font_badge)

        # 使用者 Prompt 對話氣泡 (Top Box)
        draw.rounded_rectangle([(20, 56), (800, 120)], radius=6, fill=C_IDE_SIDEBAR, outline=(71, 85, 105), width=1)
        draw.text((32, 64), "USER PROMPT / CLINICAL INTENT:", fill=C_ACCENT_BLUE, font=font_badge)
        draw.text((32, 84), m["user_prompt"], fill=C_TEXT_WHITE, font=font_sub)

        # Antigravity 執行與工具調用區 (Middle Box)
        draw.rounded_rectangle([(20, 132), (800, 420)], radius=6, fill=(10, 15, 30), outline=C_BORDER, width=1)
        draw.text((32, 142), m["action_title"], fill=C_ACCENT_TEAL, font=font_badge)

        curr_y = 170
        for step in m["steps"]:
            draw.text((32, curr_y), step, fill=C_TEXT_MUTED if not step.startswith("✓") and not step.startswith("🚨") else C_TEXT_WHITE, font=font_code)
            curr_y += 54

        # 底部狀態列 (Bottom Status Bar)
        draw.rectangle([(0, 440), (819, 479)], fill=(20, 30, 48))
        draw.line([(0, 440), (819, 440)], fill=C_BORDER, width=1)
        draw.text((20, 450), m["status"], fill=m["badge_col"], font=font_badge)

        out_path = os.path.join(img_dir, f"{m['id']}.png")
        img.convert("RGB").save(out_path, quality=95)
        generated_paths[m["id"]] = out_path

    return generated_paths

# 2. 建立整合 Antigravity 介面截圖的 12 頁極致簡報
def build_screenshot_integrated_presentation():
    ui_images = create_antigravity_ui_mockups()
    cover_bg, inner_bg = generate_notebooklm_assets()
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # NotebookLM 色彩
    C_CARD_BG = RGBColor(255, 255, 255)
    C_BORDER = RGBColor(226, 232, 240)
    C_TITLE = RGBColor(15, 23, 42)
    C_BODY = RGBColor(51, 65, 85)
    C_MUTED = RGBColor(100, 116, 139)
    C_BLUE = RGBColor(2, 132, 199)
    C_TEAL = RGBColor(13, 148, 136)
    C_CORAL = RGBColor(225, 29, 72)
    C_AMBER = RGBColor(217, 119, 6)
    C_PURPLE = RGBColor(124, 58, 237)

    def add_base_slide(title_text, subtitle_text=None, category="ANTIGRAVITY VIBE CODING"):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(inner_bg, 0, 0, prs.slide_width, prs.slide_height)

        # 頂部 NotebookLM 膠囊標籤
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.36), Inches(3.6), Inches(0.32))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(241, 245, 249)
        badge.line.color.rgb = RGBColor(203, 213, 225)
        badge.line.width = Pt(0.75)

        tb_b = slide.shapes.add_textbox(Inches(0.8), Inches(0.36), Inches(3.6), Inches(0.32))
        tf_b = tb_b.text_frame
        tf_b.margin_left = Inches(0.12)
        tf_b.margin_top = Inches(0.04)
        p_b = tf_b.paragraphs[0]
        p_b.text = f"✦ {category}"
        p_b.font.name = "Segoe UI"
        p_b.font.size = Pt(9.5)
        p_b.font.bold = True
        p_b.font.color.rgb = C_BLUE

        # 標題文字方塊
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.74), Inches(11.7), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Microsoft JhengHei"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = C_TITLE

        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = "Microsoft JhengHei"
            p2.font.size = Pt(11.5)
            p2.font.color.rgb = C_MUTED
            p2.space_before = Pt(3)

        return slide

    def add_card(slide, left, top, width, height, title, content_list, badge=None, bg_color=C_CARD_BG, border_color=C_BORDER, header_color=C_BLUE):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.25)

        tb = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.18), width - Inches(0.44), height - Inches(0.36))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Microsoft JhengHei"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = header_color

        if badge:
            p.text = f"{title}  ·  {badge}"

        for item in content_list:
            p_item = tf.add_paragraph()
            p_item.text = item
            p_item.font.name = "Microsoft JhengHei"
            p_item.font.size = Pt(10.5)
            p_item.font.color.rgb = C_BODY
            p_item.space_before = Pt(4.0)
            if item.startswith("•") or item.startswith("👉") or item.startswith("✓") or item.startswith("🚨") or item.startswith("【") or item.startswith("1.") or item.startswith("2.") or item.startswith("3.") or item.startswith("4.") or item.startswith("🔍") or item.startswith("💬") or item.startswith("🛠️"):
                p_item.font.bold = False
                p_item.font.color.rgb = C_TITLE

    def add_split_slide(slide_num, title, subtitle, category, card_title, content_list, img_key, header_color=C_BLUE, badge_text="CASE STUDY"):
        slide = add_base_slide(f"{slide_num}. {title}", subtitle, category=category)
        
        # 左側：文字解說卡片 (寬 5.4 吋)
        add_card(
            slide, Inches(0.8), Inches(1.85), Inches(5.4), Inches(5.1),
            card_title, content_list, badge=badge_text, header_color=header_color
        )
        
        # 右側：Antigravity 介面截圖 Mockup (寬 6.0 吋)
        img_path = ui_images[img_key]
        slide.shapes.add_picture(img_path, Inches(6.5), Inches(1.85), Inches(6.0), Inches(5.1))
        return slide

    # =========================================================================
    # P1: 封面
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    slide1.shapes.add_picture(cover_bg, 0, 0, prs.slide_width, prs.slide_height)

    tag_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.3), Inches(4.2), Inches(0.36))
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = RGBColor(238, 242, 255)
    tag_box.line.color.rgb = RGBColor(199, 210, 254)
    tag_box.line.width = Pt(1)
    
    t_tag = tag_box.text_frame
    p_t = t_tag.paragraphs[0]
    p_t.text = "✦ ANTIGRAVITY AGENTIC AI CASE STUDY"
    p_t.font.name = "Segoe UI"
    p_t.font.size = Pt(10)
    p_t.font.bold = True
    p_t.font.color.rgb = C_BLUE

    tb1 = slide1.shapes.add_textbox(Inches(1.2), Inches(1.85), Inches(10.5), Inches(4.8))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "何謂 Vibe Coding？"
    p1.font.name = "Microsoft JhengHei"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = C_TITLE

    p2 = tf1.add_paragraph()
    p2.text = "以 Google Antigravity 為智能副駕：打造「台灣預防接種指南助手」全景實戰"
    p2.font.name = "Microsoft JhengHei"
    p2.font.size = Pt(21)
    p2.font.bold = True
    p2.font.color.rgb = C_BLUE
    p2.space_before = Pt(10)

    p3 = tf1.add_paragraph()
    p3.text = "全篇圖解：每頁附 Antigravity IDE 實機介面 ‧ 意圖引導 ‧ 工具調用 ‧ Rust/Wasm 跨端發布"
    p3.font.name = "Microsoft JhengHei"
    p3.font.size = Pt(13.5)
    p3.font.color.rgb = C_MUTED
    p3.space_before = Pt(16)

    p4 = tf1.add_paragraph()
    p4.text = "桃園 吳鎮宇親子耳鼻喉科診所 ‧ 臨床系統真實開發歷程完整剖析"
    p4.font.name = "Microsoft JhengHei"
    p4.font.size = Pt(11.5)
    p4.font.color.rgb = C_TEAL
    p4.space_before = Pt(8)

    # =========================================================================
    # P2: Vibe Coding 核心概念 (圖文對照)
    # =========================================================================
    add_split_slide(
        1, "什麼是 Vibe Coding？從「語法撰寫」到「意圖指揮」",
        "開發者不再手敲代碼，而是透過人機對話引導 Antigravity 生成完整功能",
        "01 · CORE DEFINITION",
        "💡 核心心智模型轉變",
        [
            "• 【傳統開發】：人類是大腦也是打字員，需要記憶語法、翻閱文件、逐行 Debug。",
            "• 【Vibe Coding】：人類是【總架構師與產品總監】，AI 是【全天候資深全端工程師】。",
            "• 人類專注於「臨床需求、使用者體驗、醫療法規正確性 (The Vibe)」；AI 負責底層跨平台代碼與打包。",
            "",
            "👉 『你只需看到問題、提出想法、驗證結果。代碼自然會到位。』"
        ],
        "ui_p2_intent",
        header_color=C_BLUE,
        badge_text="Mindset"
    )

    # =========================================================================
    # P3: Antigravity 工具鏈體系 (圖文對照)
    # =========================================================================
    add_split_slide(
        2, "Antigravity 如何驅動 Vibe Coding？Agentic 核心工具機制",
        "超越一般 Chatbot——Antigravity 具備終端執行、精確代碼修補與非同步任務調度能力",
        "02 · ANTIGRAVITY ENGINE",
        "🛠️ 關鍵工具體系 (Tool Calling)",
        [
            "• **`run_command`**：直接在 PowerShell 執行編譯、Git 推送與 APK 簽署。",
            "• **`replace_file_content`**：對數千行代碼進行精準行級局部替換，絕不粗暴覆蓋。",
            "• **`search_web`**：自主聯網爬取最新官方醫學公報與法規。",
            "• **`invoke_subagent`**：派發背景子代理平行執行深色主題重構。",
            "",
            "✓ 全自動閉環：AI 自動編譯並直接讀取終端機輸出排錯。"
        ],
        "ui_p3_tools",
        header_color=C_TEAL,
        badge_text="Tool Pipeline"
    )

    # =========================================================================
    # P4: 實戰一 - 複雜醫療法規轉化 (圖文對照)
    # =========================================================================
    add_split_slide(
        3, "實戰一：自然語言 ➔ 嚴謹醫學算法 (ACIP 遲打補打規則)",
        "Antigravity 如何將衛福部疾管署厚達數十頁的時程規範，轉化為零失誤的 Rust 演算法",
        "03 · MEDICAL LOGIC",
        "📋 複雜臨床規則轉化",
        [
            "• 五合一第 4 劑需滿 1 歲 6 個月，且與第 3 劑至少隔 6 個月 (180天)。",
            "• PCV13 若在 7-11 個月補打第 1 劑，總劑次將由 4 劑縮減為 3 劑。",
            "• 活性減毒疫苗（水痘 vs MMR）同天施打可，若不同天需間隔至少 28 天。",
            "",
            "✓ 在 `vaccine-core` 建立強型別狀態機，毫秒級計算下次最短安全日！"
        ],
        "ui_p4_medical",
        header_color=C_CORAL,
        badge_text="Rust State Machine"
    )

    # =========================================================================
    # P5: 實戰二 - AI 聯網查核年份出處 (圖文對照)
    # =========================================================================
    add_split_slide(
        4, "實戰二：AI 自主法規查核與文獻引用 (0~18歲生長曲線)",
        "面對使用者質疑「資料太舊？」，Antigravity 如何調用 search_web 自主求證國健署公報",
        "04 · AUTONOMOUS RESEARCH",
        "🔍 聯網求證與出處標註",
        [
            "💬 醫師質疑：「超過 7 歲生長曲線是哪一年資料？有更新的嗎？」",
            "",
            "【Antigravity 的求證行動】：",
            "• 調用 `search_web` 搜尋衛生福利部國民健康署最新公告。",
            "• 查證結論：國健署 2009 銜接標準、2010 衛署授升字第0990700680號公告，至今 2024 年最新《兒童健康手冊》仍以此為法定標準。",
            "• 程式頂部明確列出官方出處與年份，建立臨床權威信度。"
        ],
        "ui_p5_research",
        header_color=C_AMBER,
        badge_text="Web Grounding"
    )

    # =========================================================================
    # P6: 實戰三 - Wasm 零伺服器架構 (圖文對照)
    # =========================================================================
    add_split_slide(
        5, "實戰三：零後端純前端架構 (WebAssembly 跨語言通信)",
        "兼顧「診所病患隱私安全 (零外傳伺服器)」與「極速計算性能」的架構決策",
        "05 · WASM INTEGRATION",
        "🔒 隱私安全與傳輸規範",
        [
            "• 傳統 Web 需架設後端 API，病人生日傳上雲端存在資安洩漏風險。",
            "• **Antigravity 架構決策**：將 Rust 編譯為純二進位 Wasm，在瀏覽器端本地離線執行！",
            "",
            "【跨語言避坑合約】：",
            "• Rust 端使用 `serde_json::to_string` 序列化為 String。",
            "• JS 端 `fallbackInvoke` 統一 `JSON.parse` 還原物件，完美對齊變數名稱。"
        ],
        "ui_p6_wasm",
        header_color=C_PURPLE,
        badge_text="Zero Server"
    )

    # =========================================================================
    # P7: 實戰四 - Android NDK 自主排錯 (圖文對照)
    # =========================================================================
    add_split_slide(
        6, "實戰四：跨平台編譯地獄與 Antigravity 自主修復 (Self-Healing)",
        "面對 Android NDK、Gradle、Clang 連結器報錯，Antigravity 如何自主排錯並簽署 APK",
        "06 · DEVOPS & NDK",
        "🤖 自動化編譯與修復管線",
        [
            "• 將網頁封裝為 Android 原生 App 時常遇 Clang Linker 255 錯誤代碼與符號衝突。",
            "",
            "【Antigravity 自主排錯行動】：",
            "1. 自動將 `vaccine-core` crate-type 改為純 `rlib` 解除衝突。",
            "2. 配置 Clang 34 與 NDK 環境變數，自動調用 `cargo build`。",
            "3. 自動將 `.so` 複製至 Gradle `jniLibs`，並調用 `apksigner` 完成簽署。",
            "✓ 產出可直接於實機安裝之 `台灣疫苗指南助手.apk`。"
        ],
        "ui_p7_ndk",
        header_color=C_CORAL,
        badge_text="Self-Healing"
    )

    # =========================================================================
    # P8: 實戰五 - 雙生態行事曆 (圖文對照)
    # =========================================================================
    add_split_slide(
        7, "實戰五：雙生態使用者體驗 (Google 日曆 + iOS 原生相機直掃)",
        "如何用一行對話，同時滿足 Android、Windows、Mac 與 iPhone 家長的使用習慣",
        "07 · USER EXPERIENCE",
        "📱 雙生態跨平台連動",
        [
            "💬 醫師需求：「生成行事曆 QR Code，有可能增加 iOS 日曆嗎？」",
            "",
            "【Antigravity 實現跨生態相容】：",
            "• 🌐 **Google 日曆**：生成網頁版加入行程連結。",
            "• 🍏 **Apple 日曆一鍵下載**：生成 RFC-5545 iCalendar (`.ics`) 檔案。",
            "• 📷 **iOS 相機直掃 QR Code**：切換為 `BEGIN:VEVENT` 條碼，iPhone 內建相機一照直接彈出加入行程按鈕！"
        ],
        "ui_p8_ical",
        header_color=C_BLUE,
        badge_text="Dual Calendar"
    )

    # =========================================================================
    # P9: 實戰六 - 臨床細節與安全防呆 (圖文對照)
    # =========================================================================
    add_split_slide(
        8, "實戰六：臨床安全防呆機制 (自訂預約日與提早醫學理由)",
        "從「死板的試算表」進化為「有醫學智慧的對話式互動輔助系統」",
        "08 · SAFETY & LOGIC",
        "🩺 臨床動態安全評估",
        [
            "• 家長可能因工作忙碌需改期，但在彈窗修改日期時需防止提早施打。",
            "",
            "【Antigravity 植入動態醫療評估】：",
            "🚨 **提早接種（紅色警示）**：『不建議提早接種（提早了 X 天）。若早於最小月齡施打，因母體抗體干擾將導致抗體效價不足，視為無效劑次需重打！』",
            "ℹ️ **延後接種（黃色提示）**：『順延接種即可，不需從頭重打。』",
            "✓ 日期一改，Google 網址與 iOS QR Code 瞬間同步重繪！"
        ],
        "ui_p9_warning",
        header_color=C_CORAL,
        badge_text="Clinical Safety"
    )

    # =========================================================================
    # P10: 實戰七 - Antigravity Skills 專案記憶大腦 (圖文對照)
    # =========================================================================
    add_split_slide(
        9, "實戰七：知識沉澱機制 (Antigravity Skills 專案記憶大腦)",
        "解決 AI 長對話失憶痛點——透過 SKILL.md 將架構規則與指令沉澱為永久大腦",
        "09 · ANTIGRAVITY SKILLS",
        "🧠 專案大腦記憶體系",
        [
            "• 大型專案迭代數週後，對話上下文長度會耗盡。",
            "",
            "【解法：`.agents/skills/vaccine-dev/SKILL.md`】：",
            "• 記錄 Android NDK / Clang 34 編譯環境變數與路徑。",
            "• 記錄 Rust Wasm JSON 傳輸合約與變數映射表。",
            "• 記錄 0-18 歲生長曲線國健署法定年份標準。",
            "• 記錄 GitHub Pages 部署與快取強刷 SOP。",
            "✓ Antigravity 每次被喚醒時自動讀取，立即進入頂尖狀態！"
        ],
        "ui_p10_skills",
        header_color=C_BLUE,
        badge_text="Persistent Skill"
    )

    # =========================================================================
    # P11: 實戰八 - 專屬子代理平行作業 (圖文對照)
    # =========================================================================
    add_split_slide(
        10, "實戰八：平行分工與專屬子代理 (Antigravity Subagents)",
        "主代理負責架構排程，子代理負責專注重構——人機協同的大規模並行生產力",
        "10 · SUBAGENT COLLABORATION",
        "👥 子代理並行重構",
        [
            "• 面對全系統色彩重構（改為護眼大地色系），Antigravity 調用 `invoke_subagent` 派發專屬工人：",
            "  - **Role**：`CSS Eye-Friendly Reskin`",
            "  - **Mission**：專注重構 `styles.css`，更新 20+ 組色彩變數。",
            "",
            "✓ 主代理維持高層對話，子代理在背景獨立編譯並回報成果。",
            "✓ 真正實現「一人 + 多智能代理團隊」的超級個體開發模式！"
        ],
        "ui_p11_subagents",
        header_color=C_PURPLE,
        badge_text="Subagent Delegation"
    )

    # =========================================================================
    # P12: 總結與未來展望
    # =========================================================================
    slide12 = add_base_slide(
        "11. 總結：軟體開發的文藝復興時代 (Renaissance of Creation)",
        "當 Antigravity 承擔了所有工程重擔，「領域專業知識」成為最高維度的生產力",
        category="11 · CONCLUSION & FUTURE"
    )
    add_card(
        slide12, Inches(0.8), Inches(1.85), Inches(5.6), Inches(5.1),
        "🚀 本專案帶給我們的核心啟示",
        [
            "1. **領域專家就是最佳架構師**：醫師最懂臨床流程，在 Vibe Coding 下，醫師能直接主導軟體進化，省去與外包工程師數月的溝通代溝。",
            "2. **極致敏捷的創新閉環**：一個醫療點子（如自訂日期警示）從提出、法規檢驗、編譯修復到 Android/Web 發布，僅耗時數分鐘。",
            "3. **架構思維 > 語法記憶**：未來的核心競爭力是「清晰的需求定義能力」與「臨床批判性審查」。"
        ],
        badge="Key Takeaways",
        header_color=C_BLUE
    )
    add_card(
        slide12, Inches(6.8), Inches(1.85), Inches(5.7), Inches(5.1),
        "🌟 未來展望：人人皆是數位創造者",
        [
            "• 醫護人員可以為自己的診所打造專屬臨床工具系統。",
            "• 教師可以為學生打造即時互動測驗與學習進度追蹤器。",
            "• 企業主可以在一天內驗證並交付跨平台商業系統。",
            "",
            "✨ 『在 Antigravity 與 Vibe Coding 的時代，你的想像力與專業知識，就是最強大的程式語言。』"
        ],
        badge="Empowerment",
        border_color=RGBColor(254, 215, 170),
        header_color=C_AMBER
    )

    output_path = r"E:\Vaccine\Vibe_Coding_Presentation.pptx"
    prs.save(output_path)
    print(f"Ultimate 12-slide presentation with UI mockups generated: {output_path}")

if __name__ == "__main__":
    build_screenshot_integrated_presentation()
