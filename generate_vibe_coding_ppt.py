import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 寬螢幕投影片尺寸 (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # 設計主題色系 (醫療科技質感：Deep Slate Navy & Warm Teal & Emerald)
    COLOR_BG = RGBColor(15, 23, 42)          # #0F172A 深沉科技藍底
    COLOR_CARD = RGBColor(30, 41, 59)        # #1E293B 卡片底色
    COLOR_CARD_BORDER = RGBColor(51, 65, 85) # #334155 卡片邊框
    COLOR_PRIMARY = RGBColor(56, 189, 248)   # #38BDF8 天空藍高亮
    COLOR_TEAL = RGBColor(45, 212, 191)      # #2DD4BF 醫療薄荷綠
    COLOR_AMBER = RGBColor(251, 191, 36)     # #FBBF24 亮眼警示橘黃
    COLOR_TEXT_MAIN = RGBColor(248, 250, 252)# #F8FAFC 主文字白
    COLOR_TEXT_MUTED = RGBColor(148, 163, 184)# #94A3B8 說明文字灰

    def add_base_slide(title_text, subtitle_text=None):
        slide = prs.slides.add_slide(blank_layout)
        
        # 背景底色形狀
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()

        # 頂部裝飾亮條
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.5), Inches(0.15), Inches(0.85))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_PRIMARY
        bar.line.fill.background()

        # 標題文字方塊
        tb = slide.shapes.add_textbox(Inches(1.1), Inches(0.4), Inches(11.4), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Microsoft JhengHei"
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = "Microsoft JhengHei"
            p2.font.size = Pt(13)
            p2.font.color.rgb = COLOR_TEXT_MUTED
            p2.space_before = Pt(4)

        return slide

    def add_card(slide, left, top, width, height, title, content_list, badge_text=None, border_color=COLOR_CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)

        tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), height - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Microsoft JhengHei"
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY

        if badge_text:
            p.text = f"{title}  [{badge_text}]"

        for item in content_list:
            p_item = tf.add_paragraph()
            p_item.text = item
            p_item.font.name = "Microsoft JhengHei"
            p_item.font.size = Pt(12)
            p_item.font.color.rgb = COLOR_TEXT_MUTED
            p_item.space_before = Pt(6)
            if item.startswith("•") or item.startswith("👉") or item.startswith("✓") or item.startswith("🚨"):
                p_item.font.color.rgb = COLOR_TEXT_MAIN

    # -------------------------------------------------------------
    # 投影片 1：封面 (Cover Slide)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_BG
    bg1.line.fill.background()

    tb1 = slide1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(3.8))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "何謂 Vibe Coding？"
    p1.font.name = "Microsoft JhengHei"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_PRIMARY

    p2 = tf1.add_paragraph()
    p2.text = "從「台灣預防接種指南助手」看 AI 時代人機共創的全新軟體開發典範"
    p2.font.name = "Microsoft JhengHei"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_MAIN
    p2.space_before = Pt(12)

    p3 = tf1.add_paragraph()
    p3.text = "實務專案案例剖析 ‧ 概念解析 ‧ 人機角色蛻變 ‧ 醫療輔助系統全端落地"
    p3.font.name = "Microsoft JhengHei"
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_TEAL
    p3.space_before = Pt(20)

    # -------------------------------------------------------------
    # 投影片 2：什麼是 Vibe Coding？
    # -------------------------------------------------------------
    slide2 = add_base_slide(
        "1. 什麼是 Vibe Coding？核心定義與概念演進",
        "由前 Tesla AI 總監、OpenAI 共同創辦人 Andrej Karpathy 於 2025 年初提出的顛覆性概念"
    )
    add_card(
        slide2, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2),
        "💡 Vibe Coding 的核心定義",
        [
            "• 開發者不再逐字手寫語法，而是以「直覺、意圖與自然語言對話」引導 AI。",
            "• 開發者的角色轉變為【總架構師】與【產品質檢官 (Director & Reviewer)】。",
            "• 重點放在「想要什麼體驗 (Vibe) 與核心邏輯」，讓 AI Agent 負責底層實作、修復報錯、編譯與跨平台部署。",
            "👉 『You just see stuff, say stuff, run stuff, and copy-paste stuff. It works.』"
        ],
        badge_text="全新心智模型"
    )
    add_card(
        slide2, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2),
        "⚡ 傳統開發 vs. Vibe Coding 對比",
        [
            "【傳統寫程式 (Traditional Coding)】",
            "• 手動查閱 API 文件、記憶語言語法細節",
            "• 耗費大量時間手寫刻苦的 boilerplate 樣板代碼",
            "• 遭遇跨平台 (Web/Android/NDK) 編譯地獄時耗費數天排錯",
            "",
            "【Vibe Coding 人機協作】",
            "• 領域專家（如醫師/產品經理）直接以臨床業務邏輯下達指令",
            "• AI 即刻進行法規文獻查證、Rust 算法撰寫、Wasm 封裝、APK 簽署",
            "• 敏捷迭代：一個需求在數分鐘內完成跨平台發布"
        ],
        border_color=COLOR_TEAL
    )

    # -------------------------------------------------------------
    # 投影片 3：經典案例剖析 - 台灣預防接種指南助手
    # -------------------------------------------------------------
    slide3 = add_base_slide(
        "2. 實戰案例剖析：台灣預防接種指南助手",
        "以真實複雜的醫療系統為例——如何透過 Vibe Coding 完成四端同步開發"
    )
    add_card(
        slide3, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.2),
        "🏥 臨床專業意圖 (Doctor's Intent)",
        [
            "• 「診所需要算公自費疫苗時程」",
            "• 「7歲以上兒童有生長曲線嗎？」",
            "• 「要能支援 iOS 相機直掃日曆」",
            "• 「家長改日期如果提早，要給予醫療理由警示」",
            "",
            "👉 使用者完全不需要寫任何一行 C++/Rust/Gradle，只需專注於【醫療核心邏輯】與【使用體驗】。"
        ],
        badge_text="Prompt as Logic"
    )
    add_card(
        slide3, Inches(4.8), Inches(1.6), Inches(3.7), Inches(5.2),
        "🤖 AI Agent 自主架構與執行",
        [
            "• 即時上網搜尋並查核 2024 最新衛福部國健署標準",
            "• 撰寫高效且嚴謹的 Rust 演算法核心 (vaccine-core)",
            "• 編譯 WebAssembly 驅動 GitHub Pages 純前端運算",
            "• 自動配置 Android NDK / Clang 工具鏈打包簽署 APK",
            "• 生成 iCalendar (.ics) / VEVENT 規範數據流"
        ],
        badge_text="Agent Execution",
        border_color=COLOR_PRIMARY
    )
    add_card(
        slide3, Inches(8.8), Inches(1.6), Inches(3.7), Inches(5.2),
        "🚀 跨平台成果落地 (Deliverables)",
        [
            "1. 🌐 GitHub Pages 靜態無伺服器網站",
            "2. 📱 Android 原生 arm64 簽署 APK",
            "3. 💻 Windows 單一免安裝獨立綠色版",
            "4. 🍏 iOS 相機直掃 iCal 日曆提醒",
            "",
            "✨ 從單一構想到四端落地，全由對話驅動！"
        ],
        badge_text="Multi-Platform",
        border_color=COLOR_TEAL
    )

    # -------------------------------------------------------------
    # 投影片 4：Vibe Coding 運作流程
    # -------------------------------------------------------------
    slide4 = add_base_slide(
        "3. Vibe Coding 的四步驟高頻迭代循環",
        "從自然語言對話到跨平台生產環境交付的完整閉環"
    )
    steps = [
        ("Step 1: 意圖表達 (Intent)", "使用者提出自然語言需求\n例如：「加入iOS行事曆QR Code」", COLOR_PRIMARY),
        ("Step 2: 查核與架構 (Plan)", "AI 搜尋法規、規劃資料結構\n設計前後端資料串接格式", COLOR_TEAL),
        ("Step 3: 實現與編譯 (Build)", "AI 撰寫程式碼、呼叫編譯器\n自動解決語法錯誤與依賴項", COLOR_AMBER),
        ("Step 4: 驗證與發布 (Deploy)", "產出 APK、推送到 GitHub\n使用者即刻在手機/網頁體驗", COLOR_PRIMARY),
    ]
    for i, (title, desc, color) in enumerate(steps):
        left = Inches(0.8 + i * 2.95)
        add_card(
            slide4, left, Inches(1.8), Inches(2.8), Inches(4.8),
            title,
            [desc, "", "👉 立即進入下一輪反饋循環 (Feedback Loop)"],
            border_color=color
        )

    # -------------------------------------------------------------
    # 投影片 5：關鍵支柱與成功要素
    # -------------------------------------------------------------
    slide5 = add_base_slide(
        "4. Vibe Coding 成功落地的三大支柱",
        "為什麼本專案能做到又快又穩？底層的技術架構支撐"
    )
    add_card(
        slide5, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.2),
        "1. Skill 與上下文沉澱",
        [
            "• 建立專屬 SKILL.md 開發規範",
            "• 將 NDK 環境變數、Wasm 序列化規範、Git 部署指令結構化",
            "• 讓 AI 每次重啟都能維持統一水準，不遺忘專案背景知識"
        ],
        badge_text="Persistent Memory"
    )
    add_card(
        slide5, Inches(4.8), Inches(1.6), Inches(3.7), Inches(5.2),
        "2. 工具調用能力 (Tool Calling)",
        [
            "• 具備終端機 (PowerShell)、網路搜尋、檔案精確替換工具",
            "• 遇到 Android 編譯錯誤時，AI 能自主讀取報錯並修正 Rust / Gradle 配置",
            "• 真正的【自主除錯 (Self-Healing)】"
        ],
        badge_text="Autonomous Agent",
        border_color=COLOR_PRIMARY
    )
    add_card(
        slide5, Inches(8.8), Inches(1.6), Inches(3.7), Inches(5.2),
        "3. 領域專業引導 (Domain Knowledge)",
        [
            "• 人類的不可替代性在於【臨床專業判斷】",
            "• 審核疫苗間隔邏輯是否合規、評估使用者互動體驗",
            "• 領域專家 + 強大 AI = 專業級系統迅速誕生"
        ],
        badge_text="Human-in-the-Loop",
        border_color=COLOR_TEAL
    )

    # -------------------------------------------------------------
    # 投影片 6：總結與未來展望
    # -------------------------------------------------------------
    slide6 = add_base_slide(
        "5. 總結：軟體開發的文藝復興時代",
        "每個人都能成為自己領域的數位架構師"
    )
    add_card(
        slide6, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2),
        "🚀 核心啟示 (Key Takeaways)",
        [
            "1. 程式碼不再是壁壘：思想與領域知識 (Domain Knowledge) 才是核心價值。",
            "2. 敏捷的極致體現：從醫療想法到多端上線僅需數小時，大幅降低創新試錯成本。",
            "3. 架構思維 > 語法記憶：掌握系統設計與需求拆解能力，比死記 API 更具競爭力。"
        ],
        badge_text="Summary"
    )
    add_card(
        slide6, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2),
        "🌟 未來展望 (The Future of Creation)",
        [
            "• 醫生可以為診所打造最合身的醫療輔助系統",
            "• 老師可以為學生建構客製化的互動教學軟體",
            "• 創業家可以在一天內驗證並發布 MVP 產品",
            "",
            "✨ 『在 AI 時代，你的想像力與專業知識，就是最強大的程式語言。』"
        ],
        border_color=COLOR_AMBER
    )

    output_path = r"E:\Vaccine\Vibe_Coding_實戰案例解析_台灣預防接種指南.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
