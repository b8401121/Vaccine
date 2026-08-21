import os
import math
import numpy as np
from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def generate_notebooklm_assets():
    width, height = 1920, 1080
    bg_dir = r"E:\Vaccine\ppt_assets_nlm"
    os.makedirs(bg_dir, exist_ok=True)

    # 1. 封面背景：溫潤紙質米白漸層 + 極細優雅網格 + 右側柔和彩色光暈
    cover_img = Image.new("RGBA", (width, height), (250, 249, 246, 255))
    d_cover = ImageDraw.Draw(cover_img)
    
    for y in range(height):
        r = int(250 + (244 - 250) * (y / height))
        g = int(249 + (246 - 249) * (y / height))
        b = int(246 + (240 - 246) * (y / height))
        d_cover.line([(0, y), (width, y)], fill=(r, g, b, 255))
        
    for gx in range(0, width, 48):
        for gy in range(0, height, 48):
            d_cover.ellipse([gx-1, gy-1, gx+1, gy+1], fill=(226, 232, 240, 140))

    glow_overlay = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    d_glow = ImageDraw.Draw(glow_overlay)
    
    d_glow.ellipse([width - 650, -100, width + 250, 750], fill=(219, 234, 254, 120)) # soft blue
    d_glow.ellipse([width - 450, 200, width + 350, 950], fill=(204, 251, 241, 100)) # soft teal
    d_glow.ellipse([width - 750, 450, width - 50, 1150], fill=(254, 243, 199, 90))  # soft amber

    glow_blur = glow_overlay.filter(ImageFilter.GaussianBlur(radius=80))
    cover_final = Image.alpha_composite(cover_img, glow_blur)
    
    cover_path = os.path.join(bg_dir, "nlm_cover_bg.png")
    cover_final.convert("RGB").save(cover_path, quality=95)

    # 2. 內頁背景：純淨極簡柔白 + 網格底紋
    inner_img = Image.new("RGBA", (width, height), (252, 251, 249, 255))
    d_inner = ImageDraw.Draw(inner_img)
    for y in range(height):
        r = int(253 + (248 - 253) * (y / height))
        g = int(252 + (249 - 252) * (y / height))
        b = int(250 + (246 - 250) * (y / height))
        d_inner.line([(0, y), (width, y)], fill=(r, g, b, 255))
        
    for gx in range(0, width, 48):
        for gy in range(0, height, 48):
            d_inner.ellipse([gx-1, gy-1, gx+1, gy+1], fill=(226, 232, 240, 100))

    inner_path = os.path.join(bg_dir, "nlm_inner_bg.png")
    inner_img.convert("RGB").save(inner_path, quality=95)

    return cover_path, inner_path

def build_ultimate_presentation():
    cover_bg, inner_bg = generate_notebooklm_assets()
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # NotebookLM 頂級美學調色盤
    C_CARD_BG = RGBColor(255, 255, 255)         # #FFFFFF 純白質感卡片
    C_BORDER = RGBColor(226, 232, 240)          # #E2E8F0 極細淺灰邊框
    
    C_TITLE = RGBColor(15, 23, 42)              # #0F172A 深沉石墨黑
    C_BODY = RGBColor(51, 65, 85)               # #334155 炭灰內文
    C_MUTED = RGBColor(100, 116, 139)           # #64748B 質感輔助灰
    
    C_BLUE = RGBColor(2, 132, 199)              # #0284C7 品牌湛藍
    C_TEAL = RGBColor(13, 148, 136)             # #0D9488 薄荷深青
    C_CORAL = RGBColor(225, 29, 72)             # #E11D48 亮點珊瑚紅
    C_AMBER = RGBColor(217, 119, 6)             # #D97706 琥珀棕
    C_PURPLE = RGBColor(124, 58, 237)           # #7C3AED 科技紫

    def add_base_slide(title_text, subtitle_text=None, category="ANTIGRAVITY VIBE CODING"):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(inner_bg, 0, 0, prs.slide_width, prs.slide_height)

        # 頂部 NotebookLM 膠囊標籤
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.38), Inches(3.6), Inches(0.32))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(241, 245, 249)
        badge.line.color.rgb = RGBColor(203, 213, 225)
        badge.line.width = Pt(0.75)

        tb_b = slide.shapes.add_textbox(Inches(0.8), Inches(0.38), Inches(3.6), Inches(0.32))
        tf_b = tb_b.text_frame
        tf_b.margin_left = Inches(0.12)
        tf_b.margin_top = Inches(0.04)
        p_b = tf_b.paragraphs[0]
        p_b.text = f"✦ {category}"
        p_b.font.name = "Segoe UI"
        p_b.font.size = Pt(9.5)
        p_b.font.bold = True
        p_b.font.color.rgb = C_BLUE

        # 標題文字方塊
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.76), Inches(11.7), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Microsoft JhengHei"
        p.font.size = Pt(23)
        p.font.bold = True
        p.font.color.rgb = C_TITLE

        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = "Microsoft JhengHei"
            p2.font.size = Pt(11.5)
            p2.font.color.rgb = C_MUTED
            p2.space_before = Pt(3)

        return slide

    def add_card(slide, left, top, width, height, title, content_list, badge=None, bg_color=C_CARD_BG, border_color=C_BORDER, header_color=C_BLUE):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.25)

        tb = slide.shapes.add_textbox(left + Inches(0.24), top + Inches(0.2), width - Inches(0.48), height - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Microsoft JhengHei"
        p.font.size = Pt(15.5)
        p.font.bold = True
        p.font.color.rgb = header_color

        if badge:
            p.text = f"{title}  ·  {badge}"

        for item in content_list:
            p_item = tf.add_paragraph()
            p_item.text = item
            p_item.font.name = "Microsoft JhengHei"
            p_item.font.size = Pt(10.8)
            p_item.font.color.rgb = C_BODY
            p_item.space_before = Pt(4.0)
            if item.startswith("•") or item.startswith("👉") or item.startswith("✓") or item.startswith("🚨") or item.startswith("【") or item.startswith("1.") or item.startswith("2.") or item.startswith("3.") or item.startswith("4.") or item.startswith("🔍") or item.startswith("💬") or item.startswith("🛠️"):
                p_item.font.bold = False
                p_item.font.color.rgb = C_TITLE

    # =========================================================================
    # 投影片 1：封面 (Cover Slide)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    slide1.shapes.add_picture(cover_bg, 0, 0, prs.slide_width, prs.slide_height)

    tag_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.3), Inches(4.2), Inches(0.36))
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = RGBColor(238, 242, 255)
    tag_box.line.color.rgb = RGBColor(199, 210, 254)
    tag_box.line.width = Pt(1)
    
    t_tag = tag_box.text_frame
    p_t = t_tag.paragraphs[0]
    p_t.text = "✦ ANTIGRAVITY AGENTIC AI CASE STUDY"
    p_t.font.name = "Segoe UI"
    p_t.font.size = Pt(10)
    p_t.font.bold = True
    p_t.font.color.rgb = C_BLUE

    tb1 = slide1.shapes.add_textbox(Inches(1.2), Inches(1.85), Inches(10.5), Inches(4.8))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "何謂 Vibe Coding？"
    p1.font.name = "Microsoft JhengHei"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = C_TITLE

    p2 = tf1.add_paragraph()
    p2.text = "以 Google Antigravity 為智能副駕：打造「台灣預防接種指南助手」之全景實戰"
    p2.font.name = "Microsoft JhengHei"
    p2.font.size = Pt(21)
    p2.font.bold = True
    p2.font.color.rgb = C_BLUE
    p2.space_before = Pt(10)

    p3 = tf1.add_paragraph()
    p3.text = "探索深度 Agentic 協同：Tool Calling 工具調用 ‧ Rust/Wasm 算法 ‧ Android NDK 自動化 ‧ 醫學法規查證"
    p3.font.name = "Microsoft JhengHei"
    p3.font.size = Pt(13.5)
    p3.font.color.rgb = C_MUTED
    p3.space_before = Pt(16)

    p4 = tf1.add_paragraph()
    p4.text = "桃園 吳鎮宇親子耳鼻喉科診所 ‧ 臨床系統真實開發歷程完整拆解"
    p4.font.name = "Microsoft JhengHei"
    p4.font.size = Pt(11.5)
    p4.font.color.rgb = C_TEAL
    p4.space_before = Pt(8)

    # =========================================================================
    # 投影片 2：Vibe Coding 核心概念與本質
    # =========================================================================
    slide2 = add_base_slide(
        "1. 什麼是 Vibe Coding？從「語法撰寫」到「意圖指揮」",
        "由 Andrej Karpathy 於 2025 年提出：開發者不再手敲代碼，而是透過人機對話引導系統生長",
        category="01 · CORE DEFINITION"
    )
    add_card(
        slide2, Inches(0.8), Inches(1.85), Inches(5.6), Inches(5.1),
        "💡 核心心智模型轉變",
        [
            "• 【傳統開發】：人類是大腦也是打字員，需要記憶語法、翻閱文件、逐行 Debug。",
            "• 【Vibe Coding】：人類是【總架構師與產品總監】，AI 是【全天候資深全端工程師】。",
            "• 人類專注於「臨床需求、使用者體驗、醫療法規正確性 (The Vibe)」；AI 負責底層跨平台代碼、依賴管理與打包。",
            "",
            "👉 『你只需看到問題、提出想法、驗證結果。代碼自然會到位。』"
        ],
        badge="Mindset",
        header_color=C_BLUE
    )
    add_card(
        slide2, Inches(6.8), Inches(1.85), Inches(5.7), Inches(5.1),
        "💉 本專案對應體現：疫苗指南開發",
        [
            "【真實對話範例】",
            "💬 醫師提問：「病人不見得當時有空，可以修改日期嗎？如果提早要顯示醫療理由」",
            "",
            "【AI 自動實現】",
            "✓ 零手寫代碼：AI 自動在時間軸彈窗加入 `<input type='date'>` 選擇器",
            "✓ 內建醫學規則：自動計算提早天數，引用 CDC/ACIP 最小月齡抗體生成失效原理",
            "✓ 雙向 QR Code 同步：Google 日曆與 iOS ics 條碼即時重新渲染"
        ],
        badge="Vaccine Case",
        border_color=RGBColor(186, 230, 253),
        header_color=C_TEAL
    )

    # =========================================================================
    # 投影片 3：Antigravity Agentic 核心機制解析 (NEW)
    # =========================================================================
    slide3 = add_base_slide(
        "2. Antigravity 如何驅動 Vibe Coding？Agentic 核心機制",
        "超越普通聊天機器人——Antigravity 具備終端執行、檔案操作、子代理與排程能力的自主代理體系",
        category="02 · ANTIGRAVITY ENGINE"
    )
    add_card(
        slide3, Inches(0.8), Inches(1.85), Inches(5.6), Inches(5.1),
        "🛠️ Antigravity 的關鍵工具體系 (Tool Calling)",
        [
            "• **`run_command`**：直接在 Windows PowerShell 執行編譯、Git 推送與 APK 簽署。",
            "• **`replace_file_content`**：對數千行代碼進行精準行級局部替換，絕不粗暴覆蓋。",
            "• **`search_web` / `read_url_content`**：自主聯網爬取最新官方醫學公報與法規。",
            "• **`invoke_subagent`**：派發背景子代理平行執行深色主題重構或大範圍重構。"
        ],
        badge="Autonomous Tools",
        header_color=C_BLUE
    )
    add_card(
        slide3, Inches(6.8), Inches(1.85), Inches(5.7), Inches(5.1),
        "⚡ 為什麼比 Cursor / Copilot 更適合全端？",
        [
            "1. **全自動閉環 (End-to-End Execution)**：",
            "   不需要人類手動複製代碼到終端機貼上，Antigravity 自己編譯並看編譯輸出。",
            "2. **非同步背景任務 (Reactive Wakeup)**：",
            "   編譯 Rust .so 或 Gradle APK 需要 30 秒，AI 會進入背景監聽，完成時自動喚醒。",
            "3. **思維鏈推理 (Thought-Action Loop)**：",
            "   先思考架構 ➔ 調用工具 ➔ 驗證結果 ➔ 修正問題 ➔ 回覆使用者。"
        ],
        badge="Beyond Copilot",
        border_color=RGBColor(167, 243, 208),
        header_color=C_TEAL
    )

    # =========================================================================
    # 投影片 4：實例一 - 複雜醫療法規提煉
    # =========================================================================
    slide4 = add_base_slide(
        "3. 實戰一：自然語言 ➔ 嚴謹醫學算法 (ACIP 遲打與補打規則)",
        "Antigravity 如何將衛福部疾管署厚達數十頁的時程規範，轉化為零失誤的 Rust 演算法",
        category="03 · MEDICAL LOGIC"
    )
    add_card(
        slide4, Inches(0.8), Inches(1.85), Inches(5.6), Inches(5.1),
        "📋 複雜的醫學業務需求 (Clinical Rules)",
        [
            "• 五合一第 4 劑需滿 1 歲 6 個月，且與第 3 劑至少隔 6 個月 (180天)。",
            "• 13價肺炎鏈球菌 (PCV13) 若在 7-11 個月補打第 1 劑，總劑次將由 4 劑縮減為 3 劑。",
            "• 活性減毒疫苗（水痘 vs MMR）同天施打可，若不同天需間隔至少 28 天。",
            "",
            "👉 傳統手寫需數週梳理龐大巢狀 if-else 邏輯。"
        ],
        badge="Clinical Complexity",
        header_color=C_CORAL
    )
    add_card(
        slide4, Inches(6.8), Inches(1.85), Inches(5.7), Inches(5.1),
        "⚙️ Antigravity 的演算法實現 (vaccine-core)",
        [
            "• 在 `vaccine-core/src/lib.rs` 撰寫強型別枚舉與日期計算狀態機：",
            "  - `calculate_catch_up()`: 輸入上次接種日與劑次，秒級回傳下次最短安全日。",
            "  - `co_admin_guide`: 自動比對同次同時接種組合與左右大腿/手臂施打部位。",
            "",
            "✓ 臨床精準性：完全吻合疾管署 ACIP 最低安全年齡限制與間隔法規。",
            "✓ 高性能：Rust 毫秒級運算，零延遲輸出。"
        ],
        badge="Rust Architecture",
        border_color=RGBColor(167, 243, 208),
        header_color=C_TEAL
    )

    # =========================================================================
    # 投影片 5：實例二 - Antigravity 聯網查核年代出處
    # =========================================================================
    slide5 = add_base_slide(
        "4. 實戰二：AI 自主法規查核與文獻引用 (0~18歲生長曲線)",
        "面對使用者質疑「資料太舊？」，Antigravity 如何調用 Search 工具自主求證國健署公報",
        category="04 · AUTONOMOUS RESEARCH"
    )
    add_card(
        slide5, Inches(0.8), Inches(1.85), Inches(5.6), Inches(5.1),
        "🔍 使用者反饋與挑戰 (User Feedback)",
        [
            "💬 醫師質疑：「超過 7 歲有資料嗎？務必確定是台灣最新資料，程式中要註明是哪一年資料。資料好像太舊了，有更新的嗎？」",
            "",
            "【Antigravity 的自主求證行動】：",
            "• 調用 `search_web` 與 `read_url_content` 搜尋衛生福利部國民健康署最新公告。",
            "• 查證結論：國健署 2009 公告銜接標準，2010 公告 7-18 歲 BMI 標準 (Chen & Chang)，至今 2024 年最新《兒童健康手冊》仍以此為唯一法定標準。"
        ],
        badge="Truth Grounding",
        header_color=C_AMBER
    )
    add_card(
        slide5, Inches(6.8), Inches(1.85), Inches(5.7), Inches(5.1),
        "📊 程式碼中的精準落地與年份標註",
        [
            "• 在 `index.html` 與評估報告卡片頂部明確列出官方權威出處：",
            "  1. 0~5歲：國健署現行最新《兒童健康手冊》(2024版) WHO 2006 標準",
            "  2. 5~7歲：國健署 2009 年公布國人銜接標準 (Chen & Chang)",
            "  3. 7~18歲：衛署授升字第0990700680號公告 BMI 建議值 (2010)",
            "",
            "✓ 計算身高 P3~P97 百分位，並判斷過輕、正常、過重、肥胖界值。",
            "✓ 建立透明、具醫學可信度的臨床專業背書。"
        ],
        badge="Citation Transparency",
        border_color=RGBColor(254, 215, 170),
        header_color=C_BLUE
    )

    # =========================================================================
    # 投影片 6：實例三 - Wasm 零伺服器架構
    # =========================================================================
    slide6 = add_base_slide(
        "5. 實戰三：零後端純前端架構 (WebAssembly 跨語言通信)",
        "如何兼顧「診所病患隱私安全 (零外傳伺服器)」與「極速計算性能」",
        category="05 · WASM INTEGRATION"
    )
    add_card(
        slide6, Inches(0.8), Inches(1.85), Inches(5.6), Inches(5.1),
        "🔒 醫療隱私與架構決策",
        [
            "• 傳統 Web 需架設後端 API 伺服器，病人生日傳上雲端存在資安洩漏風險。",
            "",
            "【Antigravity 提出的架構解法】：",
            "• 採用 **Rust + WebAssembly (Wasm)**：將醫療運算核心編譯為純二進位 Wasm 檔案，直接在使用者瀏覽器端本地執行！",
            "• 100% 離線運算，個資不離手機與電腦。"
        ],
        badge="Zero-Server Architecture",
        header_color=C_BLUE
    )
    add_card(
        slide6, Inches(6.8), Inches(1.85), Inches(5.7), Inches(5.1),
        "⚡ 跨語言避坑合約 (JSON Serialization)",
        [
            "• 【挑戰】：JS 與 Rust 型別對齊地獄（CamelCase vs snake_case、記憶體溢位）。",
            "• 【Antigravity 自動建立通信規範】：",
            "  1. Rust 端一律使用 `serde_json::to_string` 序列化為 String。",
            "  2. `app.js` 的 `fallbackInvoke` 統一 `JSON.parse` 還原物件，並完成變數名稱映射。",
            "",
            "✓ 完美兼顧 Rust 的型別安全與 JS 的畫面彈性。"
        ],
        badge="Contract Binding",
        border_color=RGBColor(186, 230, 253),
        header_color=C_PURPLE
    )

    # =========================================================================
    # 投影片 7：實例四 - Antigravity 自動編譯 Android NDK (Self-Healing)
    # =========================================================================
    slide7 = add_base_slide(
        "6. 實戰四：跨平台編譯地獄與 Antigravity 自主修復 (Self-Healing)",
        "面對 Android NDK、Gradle、Clang 連結器報錯，Antigravity 如何自主排錯並打包簽署 APK",
        category="06 · DEVOPS & NDK"
    )
    add_card(
        slide7, Inches(0.8), Inches(1.85), Inches(5.6), Inches(5.1),
        "💣 跨平台原生開發的報錯難題",
        [
            "• 將網頁封裝為 Android 原生 App 時常遇嚴重阻礙：",
            "  - 缺少 CMake、NDK 版本不相容、Clang Linker 255 錯誤代碼",
            "  - Rust crate-type `cdylib` 與 `rlib` 在交叉編譯時符號衝突",
            "  - Gradle 打包後的 `.so` 未正確放置於 `arm64-v8a` 目錄",
            "",
            "👉 傳統開發者往往要花費數天在 StackOverflow 爬文排錯。"
        ],
        badge="Build Errors",
        header_color=C_CORAL
    )
    add_card(
        slide7, Inches(6.8), Inches(1.85), Inches(5.7), Inches(5.1),
        "🤖 Antigravity 的自主排錯管線 (run_command)",
        [
            "• AI 讀取終端機標準錯誤輸出，自動採取修復行動：",
            "  1. 修改 `vaccine-core/Cargo.toml` 為純 `rlib` 解除符號重複定義。",
            "  2. 配置 Clang `aarch64-linux-android34-clang.cmd` 與環境變數。",
            "  3. 自動執行 PowerShell 腳本將 `.so` 複製至 `jniLibs`，並調用 Gradle 與 `apksigner` 自動簽署。",
            "",
            "✓ 產出可直接於實機安裝之 `E:\\台灣疫苗指南助手.apk`。"
        ],
        badge="Self-Healing Pipeline",
        border_color=RGBColor(167, 243, 208),
        header_color=C_TEAL
    )

    # =========================================================================
    # 投影片 8：實例五 - 雙生態日曆 (Google vs iOS iCal QR)
    # =========================================================================
    slide8 = add_base_slide(
        "7. 實戰五：雙生態使用者體驗 (Google 日曆 + iOS 原生相機直掃)",
        "如何用一行對話，同時滿足 Android、Windows、Mac 與 iPhone 家長的使用習慣",
        category="07 · USER EXPERIENCE"
    )
    add_card(
        slide8, Inches(0.8), Inches(1.85), Inches(5.6), Inches(5.1),
        "📱 需求演進歷程",
        [
            "• 初始階段：僅提供 Google 日曆 Web 連結。",
            "💬 醫師提出新需求：「有可能增加 iOS 行事曆嗎？要有 QR Code」",
            "",
            "【Antigravity 實現跨生態相容】：",
            "• 🍏 **Apple 日曆一鍵下載**：生成標準 RFC-5545 iCalendar (`.ics`) 檔案，iPhone/Mac 點擊直接彈出原生加入行程視窗。",
            "• 📷 **iOS 相機直掃 QR Code**：將 QR Code 切換為 `BEGIN:VEVENT` 條碼，iPhone 內建相機一照即可直接加入日曆！"
        ],
        badge="Dual Ecosystem",
        header_color=C_BLUE
    )
    add_card(
        slide8, Inches(6.8), Inches(1.85), Inches(5.7), Inches(5.1),
        "✨ 彈性與臨床防呆設計",
        [
            "• 雙頁籤切換：`🌐 Google 日曆 QR` 與 `🍏 iOS 相機直掃 QR`。",
            "• 提醒內容完整豐富：",
            "  - 事件標題：預防接種提醒 — 滿 2 個月",
            "  - 內文備註：五合一第1劑、13價結合型肺炎鏈球菌第1劑...",
            "  - 臨床叮嚀：『提醒：請攜帶兒童預防接種紀錄黃卡與健保卡』",
            "",
            "✓ 大幅提升診所回診率與家長依從性。"
        ],
        badge="Clinical Utility",
        border_color=RGBColor(254, 215, 170),
        header_color=C_AMBER
    )

    # =========================================================================
    # 投影片 9：實例六 - 臨床細節與安全防呆
    # =========================================================================
    slide9 = add_base_slide(
        "8. 實戰六：臨床安全防呆機制 (自訂預約日與提早醫學理由)",
        "從「死板的試算表」進化為「有醫學智慧的對話式互動輔助系統」",
        category="08 · SAFETY & LOGIC"
    )
    add_card(
        slide9, Inches(0.8), Inches(1.85), Inches(5.6), Inches(5.1),
        "⚠️ 臨床問題：家長不見得當天有空",
        [
            "• 系統原先固定算出基準日，但家長可能工作忙碌或出國需改期。",
            "• 如果家長自行提早打，會引發嚴重醫療問題（抗體效價不足失效）。",
            "",
            "【Antigravity 即時植入醫學智慧】：",
            "1. 允許在 Modal 內自由修改日期，QR Code 隨改隨更新。",
            "2. 自動計算偏差天數，動態觸發醫療警示與衛教指引。"
        ],
        badge="Interactive Logic",
        header_color=C_CORAL
    )
    add_card(
        slide9, Inches(6.8), Inches(1.85), Inches(5.7), Inches(5.1),
        "🩺 兩大動態醫療評估回饋",
        [
            "🚨 **提早接種（紅色警示）**：",
            "  『不建議提早接種（提早了 X 天）。依 CDC/ACIP 規範，若早於最小月齡施打，可能因母體抗體干擾導致抗體生成不足，視為無效劑次需重打！』",
            "",
            "ℹ️ **延後接種（黃色提示）**：",
            "  『延後 X 天提醒。若因發燒生病需延後，直接順延接種即可，不需從頭重打。』",
            "",
            "✓ 將小兒科醫師的臨床思考邏輯，無縫編織入軟體中。"
        ],
        badge="Clinical Feedback",
        border_color=RGBColor(254, 202, 202),
        header_color=C_CORAL
    )

    # =========================================================================
    # 投影片 10：實例七 - Antigravity Skill 知識沉澱
    # =========================================================================
    slide10 = add_base_slide(
        "9. 實戰七：知識沉澱機制 (Antigravity Skills 專案記憶大腦)",
        "解決 AI 長對話失憶痛點——透過 SKILL.md 將架構規則與指令沉澱為永久大腦",
        category="09 · ANTIGRAVITY SKILLS"
    )
    add_card(
        slide10, Inches(0.8), Inches(1.85), Inches(5.6), Inches(5.1),
        "🧠 為什麼需要 Antigravity Skill 文件？",
        [
            "• 大型專案迭代數週後，對話上下文長度會耗盡。",
            "• 傳統 Prompt 每次都要重複告知環境路徑、醫療規範，極度費時。",
            "",
            "【本專案的解法：`.agents/skills/vaccine-dev/`】：",
            "• 將 NDK 工具鏈路徑、JSON 序列化合約、生長曲線引用年份、部署流程全部寫入 `SKILL.md`。",
            "• Antigravity 每次被喚醒時，自動調用 Skill，立即進入頂尖資深工程師狀態！"
        ],
        badge="Skill System",
        header_color=C_BLUE
    )
    add_card(
        slide10, Inches(6.8), Inches(1.85), Inches(5.7), Inches(5.1),
        "📦 專案資產結構 (Project Asset Tree)",
        [
            "Vaccine/",
            "├── .agents/skills/vaccine-dev/  # Antigravity 專案技能大腦",
            "├── vaccine-core/                # Rust 核心演算法庫",
            "├── vaccine-app/                 # 前端 UI + Tauri 跨平台後端",
            "│   └── gen/android/             # Android 原生 Gradle 專案",
            "├── portable-launcher/           # Windows 獨立版啟動器",
            "└── gh-pages (branch)            # GitHub Pages 靜態發布端",
            "",
            "✓ 規範化、資產化、可持續自我演進的現代工程體系。"
        ],
        badge="Architecture Tree",
        border_color=RGBColor(167, 243, 208),
        header_color=C_TEAL
    )

    # =========================================================================
    # 投影片 11：實例八 - Antigravity Subagent 平行作業 (NEW)
    # =========================================================================
    slide11 = add_base_slide(
        "10. 實戰八：平行分工與專屬子代理 (Antigravity Subagents)",
        "主代理負責架構排程，子代理負責專注重構——人機協同的大規模並行生產力",
        category="10 · SUBAGENT COLLABORATION"
    )
    add_card(
        slide11, Inches(0.8), Inches(1.85), Inches(5.6), Inches(5.1),
        "👥 子代理架構 (Subagent Delegation)",
        [
            "• 當遇到大型耗時或專門任務（例如將整套 UI 從刺眼亮色改為護眼大地色系）：",
            "• Antigravity 可透過 `invoke_subagent` 派發專屬 Subagent：",
            "  - **Role**：`CSS Eye-Friendly Reskin`",
            "  - **Mission**：專注重構 `styles.css`，更新 20+ 組色彩變數、標籤色調與卡片陰影。",
            "",
            "👉 主代理與子代理分工合作，不汙染主對話上下文，極大化產能！"
        ],
        badge="Parallel Workers",
        header_color=C_PURPLE
    )
    add_card(
        slide11, Inches(6.8), Inches(1.85), Inches(5.7), Inches(5.1),
        "⚡ 本專案視覺調色成果 (Linen & Slate)",
        [
            "• 子代理在背景獨立完成全面視覺重構：",
            "  - 背景：`Warm Linen & Muted Slate` 低飽和溫潤配色",
            "  - 按鈕與頁籤：醫療薄荷青藍 (`#4A7C88`) 漸層",
            "  - 標籤分類：公費綠 (`#3D7A5A`)、補助藍 (`#3A5A80`)、自費棕 (`#7A6030`)",
            "",
            "✓ 一次指令，全系統數百個 UI 元件瞬間完成美學升級並自動通過編譯！"
        ],
        badge="Visual Polish",
        border_color=RGBColor(186, 230, 253),
        header_color=C_TEAL
    )

    # =========================================================================
    # 投影片 12：總結與未來展望
    # =========================================================================
    slide12 = add_base_slide(
        "11. 總結：軟體開發的文藝復興時代 (Renaissance of Creation)",
        "當 Antigravity 承擔了所有工程重擔，「領域專業知識」成為最高維度的生產力",
        category="11 · CONCLUSION & FUTURE"
    )
    add_card(
        slide12, Inches(0.8), Inches(1.85), Inches(5.6), Inches(5.1),
        "🚀 本專案帶給我們的核心啟示",
        [
            "1. **領域專家就是最佳架構師**：醫師最懂臨床流程，在 Vibe Coding 下，醫師能直接主導軟體進化，省去與外包工程師數月的溝通代溝。",
            "2. **極致敏捷的創新閉環**：一個醫療點子（如自訂日期警示）從提出、法規檢驗、編譯修復到 Android/Web 發布，僅耗時數分鐘。",
            "3. **架構思維 > 語法記憶**：未來的核心競爭力是「清晰的需求定義能力」與「臨床批判性審查」。"
        ],
        badge="Key Takeaways",
        header_color=C_BLUE
    )
    add_card(
        slide12, Inches(6.8), Inches(1.85), Inches(5.7), Inches(5.1),
        "🌟 未來展望：人人皆是數位創造者",
        [
            "• 醫護人員可以為自己的診所打造專屬臨床工具系統。",
            "• 教師可以為學生打造即時互動測驗與學習進度追蹤器。",
            "• 企業主可以在一天內驗證並交付跨平台商業系統。",
            "",
            "✨ 『在 Antigravity 與 Vibe Coding 的時代，你的想像力與專業知識，就是最強大的程式語言。』"
        ],
        badge="Empowerment",
        border_color=RGBColor(254, 215, 170),
        header_color=C_AMBER
    )

    output_path = r"E:\Vaccine\Vibe_Coding_Presentation.pptx"
    prs.save(output_path)
    print(f"Ultimate 12-slide Antigravity presentation generated successfully: {output_path}")

if __name__ == "__main__":
    build_ultimate_presentation()
