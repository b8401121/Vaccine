import os
from PIL import Image, ImageDraw, ImageFont, ImageFilter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def generate_notebooklm_assets():
    width, height = 1920, 1080
    bg_dir = r"E:\Vaccine\ppt_assets_nlm"
    os.makedirs(bg_dir, exist_ok=True)

    cover_img = Image.new("RGBA", (width, height), (250, 249, 246, 255))
    d_cover = ImageDraw.Draw(cover_img)
    for y in range(height):
        r = int(250 + (244 - 250) * (y / height))
        g = int(249 + (246 - 249) * (y / height))
        b = int(246 + (240 - 246) * (y / height))
        d_cover.line([(0, y), (width, y)], fill=(r, g, b, 255))
        
    for gx in range(0, width, 48):
        for gy in range(0, height, 48):
            d_cover.ellipse([gx-1, gy-1, gx+1, gy+1], fill=(226, 232, 240, 140))

    glow_overlay = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    d_glow = ImageDraw.Draw(glow_overlay)
    d_glow.ellipse([width - 650, -100, width + 250, 750], fill=(219, 234, 254, 120))
    d_glow.ellipse([width - 450, 200, width + 350, 950], fill=(204, 251, 241, 100))
    d_glow.ellipse([width - 750, 450, width - 50, 1150], fill=(254, 243, 199, 90))

    glow_blur = glow_overlay.filter(ImageFilter.GaussianBlur(radius=80))
    cover_final = Image.alpha_composite(cover_img, glow_blur)
    cover_path = os.path.join(bg_dir, "nlm_cover_bg.png")
    cover_final.convert("RGB").save(cover_path, quality=95)

    inner_img = Image.new("RGBA", (width, height), (252, 251, 249, 255))
    d_inner = ImageDraw.Draw(inner_img)
    for y in range(height):
        r = int(253 + (248 - 253) * (y / height))
        g = int(252 + (249 - 252) * (y / height))
        b = int(250 + (246 - 250) * (y / height))
        d_inner.line([(0, y), (width, y)], fill=(r, g, b, 255))
        
    for gx in range(0, width, 48):
        for gy in range(0, height, 48):
            d_inner.ellipse([gx-1, gy-1, gx+1, gy+1], fill=(226, 232, 240, 100))

    inner_path = os.path.join(bg_dir, "nlm_inner_bg.png")
    inner_img.convert("RGB").save(inner_path, quality=95)

    return cover_path, inner_path


# 100% 根據使用者截圖真實還原 Antigravity Desktop App (Light/Clean UI)
def create_authentic_antigravity_desktop_ui():
    img_dir = r"E:\Vaccine\ppt_assets_auth_ui"
    os.makedirs(img_dir, exist_ok=True)

    W, H = 840, 520
    
    # 截圖真實色彩 (Antigravity 2.0 Desktop Light Theme)
    BG_WINDOW = (248, 249, 250)        # #F8F9FA 極淡灰白主背景
    BG_SIDEBAR = (241, 243, 244)       # #F1F3F4 左側導覽列底色
    BG_CHAT = (255, 255, 255)          # #FFFFFF 中間對話畫布
    BG_INPUT = (241, 243, 244)         # #F1F3F4 底部輸入框
    BG_USER_BUBBLE = (241, 245, 249)   # #F1F5F9 使用者提問灰白卡片
    BG_TOOL_BOX = (248, 250, 252)      # #F8FAFC 工具調用框
    BORDER_LIGHT = (226, 232, 240)     # #E2E8F0 極細邊框線
    BORDER_SIDEBAR = (218, 220, 224)   # #DADCE0
    
    TEXT_MAIN = (32, 33, 36)           # #202124 主文字深黑
    TEXT_MUTED = (95, 99, 104)         # #5F6368 次要文字灰
    TEXT_BLUE = (26, 115, 232)         # #1A73E8 品牌藍 (Google Blue)
    TEXT_GREEN = (24, 128, 56)         # #188038 成功綠
    TEXT_AMBER = (234, 134, 0)         # #EA8600 警示橘
    TEXT_RED = (217, 48, 37)           # #D93025 警告紅

    font_path = "C:/Windows/Fonts/msjh.ttc"
    font_bold_path = "C:/Windows/Fonts/msjhbd.ttc"
    font_code_path = "C:/Windows/Fonts/consola.ttf"

    try:
        f_menu = ImageFont.truetype(font_path, 11)
        f_title = ImageFont.truetype(font_bold_path, 12)
        f_side_btn = ImageFont.truetype(font_path, 11)
        f_prompt = ImageFont.truetype(font_bold_path, 12)
        f_body = ImageFont.truetype(font_path, 11)
        f_code = ImageFont.truetype(font_code_path, 10.5)
        f_tag = ImageFont.truetype(font_bold_path, 9.5)
    except:
        f_menu = f_title = f_side_btn = f_prompt = f_body = f_code = f_tag = ImageFont.load_default()

    data_list = [
        {
            "id": "ui_p2_intent",
            "active_topic": "Vaccine / Adding Calendar Warning",
            "user_prompt": "病人不見得當時有空，可以修改日期嗎？如果修改的日期是不建議的，如提早，顯示理由",
            "agent_thought": "Thinking Process: 1. Modify calendar modal with <input type='date'>. 2. Calculate interval diff. 3. Add CDC/ACIP maternal antibody explanation.",
            "tool_action": "Editing file: e:\\Vaccine\\vaccine-app\\src\\main.js",
            "diff_summary": "1 file changed +68 -8  ·  replace_file_content",
            "agent_reply": "已為您加入自訂日期選擇器！若家長提早施打，會跳出紅色警示說明『最小月齡限制與母體抗體干擾導致失效』；若延後則提醒『順延接種即可』，QR Code 隨改隨更新。",
            "badge": "Agent Mode"
        },
        {
            "id": "ui_p3_tools",
            "active_topic": "Vaccine / Tool Calling Engine",
            "user_prompt": "同時執行法規搜尋、程式修補、終端編譯與背景子代理",
            "agent_thought": "Orchestrating tool pipeline concurrently across 4 subsystems...",
            "tool_action": "Running command: cargo build --target aarch64-linux-android --release",
            "diff_summary": "4 tools executed  ·  Exit Code 0  ·  Duration 14.2s",
            "agent_reply": "Antigravity 已自動調用 search_web 搜尋疾管署規範、replace_file_content 局部修改代碼、run_command 在終端機編譯，並派發 Subagent 平行作業。",
            "badge": "Tool Chain"
        },
        {
            "id": "ui_p4_medical",
            "active_topic": "Vaccine / Catch-Up Algorithm",
            "user_prompt": "五合一、PCV13 遲打補打時程規則，請轉化為嚴謹的演算法",
            "agent_thought": "Translating Taiwan CDC & ACIP complex clinical schedules into Rust state machine...",
            "tool_action": "Editing file: e:\\Vaccine\\vaccine-core\\src\\lib.rs",
            "diff_summary": "1 file changed +142 -12  ·  vaccine-core (Rust)",
            "agent_reply": "已在 vaccine-core 實作狀態機：五合一第4劑自動鎖定間隔180天且滿18個月；PCV13於7~11個月補打時自動縮減為3劑(2+1)，毫秒級精準推算。",
            "badge": "Rust Algorithm"
        },
        {
            "id": "ui_p5_research",
            "active_topic": "Vaccine / MOHW Citation Grounding",
            "user_prompt": "生長曲線超過7歲是哪一年資料？有更新的嗎？務必確定是台灣最新資料",
            "agent_thought": "Searching MOHW Health Promotion Administration official datasets...",
            "tool_action": "Searching the web: 'Taiwan CDC MOHW child growth curve 0-18 BMI standard'",
            "diff_summary": "Verified 3 official legal sources  ·  MOHW (2006/2009/2010)",
            "agent_reply": "已求證國健署最新公告：現行2024《兒童健康手冊》0-5歲為WHO 2006標準，5-7歲為2009銜接標準，7-18歲為2010年衛署授升字第0990700680號公告BMI標準，已於介面完整標註年份出處。",
            "badge": "search_web"
        },
        {
            "id": "ui_p6_wasm",
            "active_topic": "Vaccine / Zero-Server Wasm Bridge",
            "user_prompt": "病人生日與健保個資不能上傳雲端，如何實現純前端離線運算？",
            "agent_thought": "Structuring WebAssembly zero-server IPC bridge with JSON serialization...",
            "tool_action": "Editing file: e:\\Vaccine\\app.js & vaccine-core/Cargo.toml",
            "diff_summary": "1 file changed +52 -4  ·  Wasm JSON contract",
            "agent_reply": "已架構 Rust -> WebAssembly 本地運算管道。病人生日與時程推算100%在瀏覽器記憶體內執行，零外傳伺服器，完全符合醫療隱私法規。",
            "badge": "WebAssembly"
        },
        {
            "id": "ui_p7_ndk",
            "active_topic": "Vaccine / Android NDK Pipeline",
            "user_prompt": "編譯 Android aarch64 出現 Clang Linker 255 錯誤與符號重複",
            "agent_thought": "Inspecting stderr: cdylib symbol conflict detected; switching to rlib...",
            "tool_action": "Running command: .\\gradlew.bat assembleArm64Release & apksigner",
            "diff_summary": "Task :app:assembleArm64Release -> SUCCESS in 10s",
            "agent_reply": "Antigravity 自動將 Cargo.toml crate-type 改為 rlib 解除符號衝突，設定 Clang 34 工具鏈，將編譯好的 .so 複製至 jniLibs 並完成 APK 簽署 (E:\\台灣疫苗指南助手.apk)。",
            "badge": "Self-Healing"
        },
        {
            "id": "ui_p8_ical",
            "active_topic": "Vaccine / iOS Calendar & QR Code",
            "user_prompt": "生成專屬行事曆連結與 QR Code，有可能增加 iOS 行事曆嗎？",
            "agent_thought": "Implementing RFC-5545 iCalendar (.ics) blob & VEVENT QR Code stream...",
            "tool_action": "Editing file: e:\\Vaccine\\vaccine-app\\src\\main.js",
            "diff_summary": "1 file changed +53 -8  ·  iOS VEVENT generator",
            "agent_reply": "已實作雙切換頁籤：『Google 日曆 QR』與『iOS 相機直掃 QR』。iPhone 內建相機直接朝向螢幕即可一鍵加入 iOS 行事曆！",
            "badge": "iOS iCal Bridge"
        },
        {
            "id": "ui_p9_warning",
            "active_topic": "Vaccine / Clinical Safety Reasoning",
            "user_prompt": "家長修改預約日期，若提早要有醫學警示，若延後要給予衛教",
            "agent_thought": "Binding real-time onchange event with medical safety interval evaluation...",
            "tool_action": "Editing file: e:\\Vaccine\\vaccine-app\\src\\index.html",
            "diff_summary": "1 file changed +162 -49  ·  Medical warning engine",
            "agent_reply": "已實作動態警示：提早施打顯示紅色警告說明『母體抗體干擾導致無效』；延後施打顯示黃色提示『直接順延即可』，日曆網址與 QR Code 隨選隨改。",
            "badge": "Clinical Safety"
        },
        {
            "id": "ui_p10_skills",
            "active_topic": "Vaccine / Persistent Skill Brain",
            "user_prompt": "更新此專案 skill",
            "agent_thought": "Updating long-term memory at .agents/skills/vaccine-dev/SKILL.md...",
            "tool_action": "Editing file: e:\\Vaccine\\.agents\\skills\\vaccine-dev\\SKILL.md",
            "diff_summary": "1 file changed +84 -30  ·  Dev Skill Architecture",
            "agent_reply": "已將 NDK 編譯變數、Wasm 序列化規範、國健署 2009/2010 法定年份與 GitHub Pages 部署 SOP 完整記錄於 SKILL.md，專案重啟時自動繼承記憶。",
            "badge": "SKILL.md"
        },
        {
            "id": "ui_p11_subagents",
            "active_topic": "Vaccine / Background Subagents",
            "user_prompt": "將全系統 UI 配色改為護眼溫潤的大地亞麻色系 (Warm Linen)",
            "agent_thought": "Spawning background subagent 'CSS Eye-Friendly Reskin' for parallel refactoring...",
            "tool_action": "Invoking Subagent: 'CSS Eye-Friendly Reskin'",
            "diff_summary": "Subagent finished in 24s  ·  styles.css refactored",
            "agent_reply": "專屬子代理已在背景獨立重構 styles.css，全面套用低飽和亞麻色與青藍色，主對話上下文零干擾，自動通過編譯！",
            "badge": "invoke_subagent"
        }
    ]

    generated_paths = {}

    for d in data_list:
        img = Image.new("RGBA", (W, H), BG_WINDOW)
        draw = ImageDraw.Draw(img)

        # 1. 最頂部 App 視窗標題列 (38px)
        draw.rectangle([(0, 0), (W-1, 38)], fill=BG_WINDOW)
        draw.line([(0, 38), (W-1, 38)], fill=BORDER_SIDEBAR, width=1)
        
        # 頂部文字選單 (Antigravity File View Window)
        draw.text((16, 12), "Antigravity", fill=TEXT_MAIN, font=f_title)
        draw.text((105, 12), "File", fill=TEXT_MUTED, font=f_menu)
        draw.text((140, 12), "View", fill=TEXT_MUTED, font=f_menu)
        draw.text((180, 12), "Window", fill=TEXT_MUTED, font=f_menu)

        # 頂部中右側工作區導覽路徑
        draw.text((250, 12), f"{d['active_topic']}", fill=TEXT_MUTED, font=f_menu)
        
        # 右上角視窗按鈕 (最小化/最大化/關閉)
        draw.text((W - 80, 10), "—   □   ✕", fill=TEXT_MUTED, font=f_menu)

        # 2. 左側導覽欄 (Left Sidebar: 190px 寬, x: 0 ~ 190)
        draw.rectangle([(0, 39), (190, H-1)], fill=BG_SIDEBAR)
        draw.line([(190, 39), (190, H-1)], fill=BORDER_SIDEBAR, width=1)

        # "+ New Conversation" 膠囊按鈕
        draw.rounded_rectangle([(12, 50), (178, 80)], radius=15, fill=(255, 255, 255), outline=BORDER_SIDEBAR, width=1)
        draw.text((26, 58), "+  New Conversation", fill=TEXT_MAIN, font=f_side_btn)

        draw.text((16, 96), "🕘  Conversation History", fill=TEXT_MUTED, font=f_side_btn)
        draw.text((16, 122), "⏱️  Scheduled Tasks", fill=TEXT_MUTED, font=f_side_btn)

        draw.line([(12, 150), (178, 150)], fill=BORDER_SIDEBAR, width=1)
        draw.text((16, 160), "Projects", fill=TEXT_MUTED, font=f_tag)

        # 專案列表
        draw.text((16, 182), "📁  Vaccine", fill=TEXT_BLUE, font=f_side_btn)
        draw.rounded_rectangle([(24, 202), (178, 226)], radius=4, fill=(230, 240, 255))
        draw.text((30, 208), "Vaccine Assistant...  now", fill=TEXT_BLUE, font=f_tag)

        draw.text((16, 236), "📁  PaxlovidWeb", fill=TEXT_MUTED, font=f_side_btn)
        draw.text((16, 260), "📁  BEER", fill=TEXT_MUTED, font=f_side_btn)
        draw.text((30, 282), "程式專案編譯 APK 諮詢", fill=TEXT_MUTED, font=f_tag)

        # 左下角 Settings
        draw.line([(12, H - 36), (178, H - 36)], fill=BORDER_SIDEBAR, width=1)
        draw.text((16, H - 24), "⚙️  Settings", fill=TEXT_MUTED, font=f_side_btn)

        # 3. 中間主對話畫布 (Chat Canvas: x: 191 ~ W-1)
        draw.rectangle([(191, 39), (W-1, H-1)], fill=BG_CHAT)

        # (A) 使用者提問氣泡 (Top Box)
        draw.rounded_rectangle([(210, 52), (W - 30, 102)], radius=8, fill=BG_USER_BUBBLE, outline=BORDER_LIGHT, width=1)
        draw.text((224, 62), "👤 User Prompt:", fill=TEXT_BLUE, font=f_tag)
        draw.text((224, 76), d["user_prompt"], fill=TEXT_MAIN, font=f_prompt)

        # (B) Antigravity Agent Thinking 區塊
        draw.text((210, 116), "🤖 Antigravity", fill=TEXT_MAIN, font=f_title)
        draw.rounded_rectangle([(W - 140, 114), (W - 30, 134)], radius=4, fill=(238, 242, 255), outline=TEXT_BLUE, width=1)
        draw.text((W - 130, 118), f"✦ {d['badge']}", fill=TEXT_BLUE, font=f_tag)

        draw.text((210, 138), f"💭 {d['agent_thought']}", fill=TEXT_MUTED, font=f_body)

        # (C) Tool Call 卡片 (中間工具執行框)
        draw.rounded_rectangle([(210, 160), (W - 30, 226)], radius=6, fill=BG_TOOL_BOX, outline=BORDER_LIGHT, width=1)
        draw.text((222, 170), f"🛠️ {d['tool_action']}", fill=TEXT_BLUE, font=f_title)
        draw.text((222, 192), f"   {d['diff_summary']}", fill=TEXT_GREEN, font=f_code)
        draw.rounded_rectangle([(W - 100, 190), (W - 42, 214)], radius=4, fill=(255, 255, 255), outline=BORDER_SIDEBAR, width=1)
        draw.text((W - 90, 196), "👁️ Review", fill=TEXT_MUTED, font=f_tag)

        # (D) Agent 最終回覆文字 (下部)
        draw.rounded_rectangle([(210, 236), (W - 30, 360)], radius=8, fill=(255, 255, 255), outline=BORDER_LIGHT, width=1)
        draw.text((224, 246), "📋 Agent Response:", fill=TEXT_GREEN, font=f_tag)
        
        # 折行繪製回覆文字
        reply_lines = []
        raw_text = d["agent_reply"]
        line_len = 38
        for i in range(0, len(raw_text), line_len):
            reply_lines.append(raw_text[i:i+line_len])
            
        cur_ry = 268
        for r_line in reply_lines:
            draw.text((224, cur_ry), r_line, fill=TEXT_MAIN, font=f_body)
            cur_ry += 22

        # (E) 底部輸入框 (Bottom Input Bar)
        draw.rounded_rectangle([(210, H - 74), (W - 30, H - 18)], radius=8, fill=BG_INPUT, outline=BORDER_SIDEBAR, width=1)
        draw.text((224, H - 64), "Ask anything, @ to mention, / for actions", fill=TEXT_MUTED, font=f_body)
        draw.text((224, H - 38), "✦  Gemini 3.7 Flash (Medium) ▾", fill=TEXT_BLUE, font=f_tag)
        draw.rounded_rectangle([(W - 60, H - 56), (W - 40, H - 36)], radius=10, fill=TEXT_BLUE)
        draw.text((W - 53, H - 53), "↑", fill=(255, 255, 255), font=f_title)

        out_path = os.path.join(img_dir, f"{d['id']}.png")
        img.convert("RGB").save(out_path, quality=95)
        generated_paths[d["id"]] = out_path

    return generated_paths


def build_final_authentic_presentation():
    ui_images = create_authentic_antigravity_desktop_ui()
    cover_bg, inner_bg = generate_notebooklm_assets()
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    C_CARD_BG = RGBColor(255, 255, 255)
    C_BORDER = RGBColor(226, 232, 240)
    C_TITLE = RGBColor(15, 23, 42)
    C_BODY = RGBColor(51, 65, 85)
    C_MUTED = RGBColor(100, 116, 139)
    C_BLUE = RGBColor(2, 132, 199)
    C_TEAL = RGBColor(13, 148, 136)
    C_CORAL = RGBColor(225, 29, 72)
    C_AMBER = RGBColor(217, 119, 6)
    C_PURPLE = RGBColor(124, 58, 237)

    def add_base_slide(title_text, subtitle_text=None, category="ANTIGRAVITY VIBE CODING"):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(inner_bg, 0, 0, prs.slide_width, prs.slide_height)

        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.36), Inches(3.6), Inches(0.32))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(241, 245, 249)
        badge.line.color.rgb = RGBColor(203, 213, 225)
        badge.line.width = Pt(0.75)

        tb_b = slide.shapes.add_textbox(Inches(0.8), Inches(0.36), Inches(3.6), Inches(0.32))
        tf_b = tb_b.text_frame
        tf_b.margin_left = Inches(0.12)
        tf_b.margin_top = Inches(0.04)
        p_b = tf_b.paragraphs[0]
        p_b.text = f"✦ {category}"
        p_b.font.name = "Segoe UI"
        p_b.font.size = Pt(9.5)
        p_b.font.bold = True
        p_b.font.color.rgb = C_BLUE

        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.74), Inches(11.7), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Microsoft JhengHei"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = C_TITLE

        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = "Microsoft JhengHei"
            p2.font.size = Pt(11.5)
            p2.font.color.rgb = C_MUTED
            p2.space_before = Pt(3)

        return slide

    def add_card(slide, left, top, width, height, title, content_list, badge=None, bg_color=C_CARD_BG, border_color=C_BORDER, header_color=C_BLUE):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.25)

        tb = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.18), width - Inches(0.44), height - Inches(0.36))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Microsoft JhengHei"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = header_color

        if badge:
            p.text = f"{title}  ·  {badge}"

        for item in content_list:
            p_item = tf.add_paragraph()
            p_item.text = item
            p_item.font.name = "Microsoft JhengHei"
            p_item.font.size = Pt(10.5)
            p_item.font.color.rgb = C_BODY
            p_item.space_before = Pt(4.0)
            if item.startswith("•") or item.startswith("👉") or item.startswith("✓") or item.startswith("🚨") or item.startswith("【") or item.startswith("1.") or item.startswith("2.") or item.startswith("3.") or item.startswith("4.") or item.startswith("🔍") or item.startswith("💬") or item.startswith("🛠️"):
                p_item.font.bold = False
                p_item.font.color.rgb = C_TITLE

    def add_split_slide(slide_num, title, subtitle, category, card_title, content_list, img_key, header_color=C_BLUE, badge_text="CASE STUDY"):
        slide = add_base_slide(f"{slide_num}. {title}", subtitle, category=category)
        add_card(
            slide, Inches(0.8), Inches(1.85), Inches(5.3), Inches(5.15),
            card_title, content_list, badge=badge_text, header_color=header_color
        )
        img_path = ui_images[img_key]
        slide.shapes.add_picture(img_path, Inches(6.4), Inches(1.85), Inches(6.1), Inches(5.15))
        return slide

    # P1: 封面
    slide1 = prs.slides.add_slide(blank_layout)
    slide1.shapes.add_picture(cover_bg, 0, 0, prs.slide_width, prs.slide_height)

    tag_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.3), Inches(4.2), Inches(0.36))
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = RGBColor(238, 242, 255)
    tag_box.line.color.rgb = RGBColor(199, 210, 254)
    tag_box.line.width = Pt(1)
    
    t_tag = tag_box.text_frame
    p_t = t_tag.paragraphs[0]
    p_t.text = "✦ ANTIGRAVITY AGENTIC AI CASE STUDY"
    p_t.font.name = "Segoe UI"
    p_t.font.size = Pt(10)
    p_t.font.bold = True
    p_t.font.color.rgb = C_BLUE

    tb1 = slide1.shapes.add_textbox(Inches(1.2), Inches(1.85), Inches(10.5), Inches(4.8))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "何謂 Vibe Coding？"
    p1.font.name = "Microsoft JhengHei"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = C_TITLE

    p2 = tf1.add_paragraph()
    p2.text = "以 Google Antigravity 為智能副駕：打造「台灣預防接種指南助手」全景實戰"
    p2.font.name = "Microsoft JhengHei"
    p2.font.size = Pt(21)
    p2.font.bold = True
    p2.font.color.rgb = C_BLUE
    p2.space_before = Pt(10)

    p3 = tf1.add_paragraph()
    p3.text = "全篇圖解：100% 還原 Antigravity 2.0 官方桌面應用實機介面 ‧ 工具調用 ‧ Rust/Wasm 跨端發布"
    p3.font.name = "Microsoft JhengHei"
    p3.font.size = Pt(13.5)
    p3.font.color.rgb = C_MUTED
    p3.space_before = Pt(16)

    p4 = tf1.add_paragraph()
    p4.text = "桃園 吳鎮宇親子耳鼻喉科診所 ‧ 臨床系統真實開發歷程完整剖析"
    p4.font.name = "Microsoft JhengHei"
    p4.font.size = Pt(11.5)
    p4.font.color.rgb = C_TEAL
    p4.space_before = Pt(8)

    # P2: Vibe Coding 核心概念
    add_split_slide(
        1, "什麼是 Vibe Coding？從「語法撰寫」到「意圖指揮」",
        "開發者不再手敲代碼，而是透過人機對話引導 Antigravity 生成完整功能",
        "01 · CORE DEFINITION",
        "💡 核心心智模型轉變",
        [
            "• 【傳統開發】：人類是大腦也是打字員，需要記憶語法、翻閱文件、逐行 Debug。",
            "• 【Vibe Coding】：人類是【總架構師與產品總監】，AI 是【全天候資深全端工程師】。",
            "• 人類專注於「臨床需求、使用者體驗、醫療法規正確性 (The Vibe)」；AI 負責底層跨平台代碼與打包。",
            "",
            "👉 『你只需看到問題、提出想法、驗證結果。代碼自然會到位。』"
        ],
        "ui_p2_intent",
        header_color=C_BLUE,
        badge_text="Mindset"
    )

    # P3: Antigravity 工具鏈體系
    add_split_slide(
        2, "Antigravity 如何驅動 Vibe Coding？Agentic 核心工具機制",
        "超越一般 Chatbot——Antigravity 具備終端執行、精確代碼修補與非同步任務調度能力",
        "02 · ANTIGRAVITY ENGINE",
        "🛠️ 關鍵工具體系 (Tool Calling)",
        [
            "• **`run_command`**：直接在 PowerShell 執行編譯、Git 推送與 APK 簽署。",
            "• **`replace_file_content`**：對數千行代碼進行精準行級局部替換，絕不粗暴覆蓋。",
            "• **`search_web`**：自主聯網爬取最新官方醫學公報與法規。",
            "• **`invoke_subagent`**：派發背景子代理平行執行深色主題重構。",
            "",
            "✓ 全自動閉環：AI 自動編譯並直接讀取終端機輸出排錯。"
        ],
        "ui_p3_tools",
        header_color=C_TEAL,
        badge_text="Tool Pipeline"
    )

    # P4: 實戰一 - 複雜醫療法規轉化
    add_split_slide(
        3, "實戰一：自然語言 ➔ 嚴謹醫學算法 (ACIP 遲打補打規則)",
        "Antigravity 如何將衛福部疾管署厚達數十頁的時程規範，轉化為零失誤的 Rust 演算法",
        "03 · MEDICAL LOGIC",
        "📋 複雜臨床規則轉化",
        [
            "• 五合一第 4 劑需滿 1 歲 6 個月，且與第 3 劑至少隔 6 個月 (180天)。",
            "• PCV13 若在 7-11 個月補打第 1 劑，總劑次將由 4 劑縮減為 3 劑。",
            "• 活性減毒疫苗（水痘 vs MMR）同天施打可，若不同天需間隔至少 28 天。",
            "",
            "✓ 在 `vaccine-core` 建立強型別狀態機，毫秒級計算下次最短安全日！"
        ],
        "ui_p4_medical",
        header_color=C_CORAL,
        badge_text="Rust State Machine"
    )

    # P5: 實戰二 - AI 聯網查核年份出處
    add_split_slide(
        4, "實戰二：AI 自主法規查核與文獻引用 (0~18歲生長曲線)",
        "面對使用者質疑「資料太舊？」，Antigravity 如何調用 search_web 自主求證國健署公報",
        "04 · AUTONOMOUS RESEARCH",
        "🔍 聯網求證與出處標註",
        [
            "💬 醫師質疑：「超過 7 歲生長曲線是哪一年資料？有更新的嗎？」",
            "",
            "【Antigravity 的求證行動】：",
            "• 調用 `search_web` 搜尋衛生福利部國民健康署最新公告。",
            "• 查證結論：國健署 2009 銜接標準、2010 衛署授升字第0990700680號公告，至今 2024 年最新《兒童健康手冊》仍以此為法定標準。",
            "• 程式頂部明確列出官方出處與年份，建立臨床權威信度。"
        ],
        "ui_p5_research",
        header_color=C_AMBER,
        badge_text="Web Grounding"
    )

    # P6: 實戰三 - Wasm 零伺服器架構
    add_split_slide(
        5, "實戰三：零後端純前端架構 (WebAssembly 跨語言通信)",
        "兼顧「診所病患隱私安全 (零外傳伺服器)」與「極速計算性能」的架構決策",
        "05 · WASM INTEGRATION",
        "🔒 隱私安全與傳輸規範",
        [
            "• 傳統 Web 需架設後端 API，病人生日傳上雲端存在資安洩漏風險。",
            "• **Antigravity 架構決策**：將 Rust 編譯為純二進位 Wasm，在瀏覽器端本地離線執行！",
            "",
            "【跨語言避坑合約】：",
            "• Rust 端使用 `serde_json::to_string` 序列化為 String。",
            "• JS 端 `fallbackInvoke` 統一 `JSON.parse` 還原物件，完美對齊變數名稱。"
        ],
        "ui_p6_wasm",
        header_color=C_PURPLE,
        badge_text="Zero Server"
    )

    # P7: 實戰四 - Android NDK 自主排錯
    add_split_slide(
        6, "實戰四：跨平台編譯地獄與 Antigravity 自主修復 (Self-Healing)",
        "面對 Android NDK、Gradle、Clang 連結器報錯，Antigravity 如何自主排錯並簽署 APK",
        "06 · DEVOPS & NDK",
        "🤖 自動化編譯與修復管線",
        [
            "• 將網頁封裝為 Android 原生 App 時常遇 Clang Linker 255 錯誤代碼與符號衝突。",
            "",
            "【Antigravity 自主排錯行動】：",
            "1. 自動將 `vaccine-core` crate-type 改為純 `rlib` 解除衝突。",
            "2. 配置 Clang 34 與 NDK 環境變數，自動調用 `cargo build`。",
            "3. 自動將 `.so` 複製至 Gradle `jniLibs`，並調用 `apksigner` 完成簽署。",
            "✓ 產出可直接於實機安裝之 `台灣疫苗指南助手.apk`。"
        ],
        "ui_p7_ndk",
        header_color=C_CORAL,
        badge_text="Self-Healing"
    )

    # P8: 實戰五 - 雙生態行事曆
    add_split_slide(
        7, "實戰五：雙生態使用者體驗 (Google 日曆 + iOS 原生相機直掃)",
        "如何用一行對話，同時滿足 Android、Windows、Mac 與 iPhone 家長的使用習慣",
        "07 · USER EXPERIENCE",
        "📱 雙生態跨平台連動",
        [
            "💬 醫師需求：「生成行事曆 QR Code，有可能增加 iOS 日曆嗎？」",
            "",
            "【Antigravity 實現跨生態相容】：",
            "• 🌐 **Google 日曆**：生成網頁版加入行程連結。",
            "• 🍏 **Apple 日曆一鍵下載**：生成 RFC-5545 iCalendar (`.ics`) 檔案。",
            "• 📷 **iOS 相機直掃 QR Code**：切換為 `BEGIN:VEVENT` 條碼，iPhone 內建相機一照直接彈出加入行程按鈕！"
        ],
        "ui_p8_ical",
        header_color=C_BLUE,
        badge_text="Dual Calendar"
    )

    # P9: 實戰六 - 臨床細節與安全防呆
    add_split_slide(
        8, "實戰六：臨床安全防呆機制 (自訂預約日與提早醫學理由)",
        "從「死板的試算表」進化為「有醫學智慧的對話式互動輔助系統」",
        "08 · SAFETY & LOGIC",
        "🩺 臨床動態安全評估",
        [
            "• 家長可能因工作忙碌需改期，但在彈窗修改日期時需防止提早施打。",
            "",
            "【Antigravity 植入動態醫療評估】：",
            "🚨 **提早接種（紅色警示）**：『不建議提早接種（提早了 X 天）。若早於最小月齡施打，因母體抗體干擾將導致抗體效價不足，視為無效劑次需重打！』",
            "ℹ️ **延後接種（黃色提示）**：『順延接種即可，不需從頭重打。』",
            "✓ 日期一改，Google 網址與 iOS QR Code 瞬間同步重繪！"
        ],
        "ui_p9_warning",
        header_color=C_CORAL,
        badge_text="Clinical Safety"
    )

    # P10: 實戰七 - Antigravity Skills 專案記憶大腦
    add_split_slide(
        9, "實戰七：知識沉澱機制 (Antigravity Skills 專案記憶大腦)",
        "解決 AI 長對話失憶痛點——透過 SKILL.md 將架構規則與指令沉澱為永久大腦",
        "09 · ANTIGRAVITY SKILLS",
        "🧠 專案大腦記憶體系",
        [
            "• 大型專案迭代數週後，對話上下文長度會耗盡。",
            "",
            "【解法：`.agents/skills/vaccine-dev/SKILL.md`】：",
            "• 記錄 Android NDK / Clang 34 編譯環境變數與路徑。",
            "• 記錄 Rust Wasm JSON 傳輸合約與變數映射表。",
            "• 記錄 0-18 歲生長曲線國健署法定年份標準。",
            "• 記錄 GitHub Pages 部署與快取強刷 SOP。",
            "✓ Antigravity 每次被喚醒時自動讀取，立即進入頂尖狀態！"
        ],
        "ui_p10_skills",
        header_color=C_BLUE,
        badge_text="Persistent Skill"
    )

    # P11: 實戰八 - 專屬子代理平行作業
    add_split_slide(
        10, "實戰八：平行分工與專屬子代理 (Antigravity Subagents)",
        "主代理負責架構排程，子代理負責專注重構——人機協同的大規模並行生產力",
        "10 · SUBAGENT COLLABORATION",
        "👥 子代理並行重構",
        [
            "• 面對全系統色彩重構（改為護眼大地色系），Antigravity 調用 `invoke_subagent` 派發專屬工人：",
            "  - **Role**：`CSS Eye-Friendly Reskin`",
            "  - **Mission**：專注重構 `styles.css`，更新 20+ 組色彩變數。",
            "",
            "✓ 主代理維持高層對話，子代理在背景獨立編譯並回報成果。",
            "✓ 真正實現「一人 + 多智能代理團隊」的超級個體開發模式！"
        ],
        "ui_p11_subagents",
        header_color=C_PURPLE,
        badge_text="Subagent Delegation"
    )

    # P12: 總結與未來展望
    slide12 = add_base_slide(
        "11. 總結：軟體開發的文藝復興時代 (Renaissance of Creation)",
        "當 Antigravity 承擔了所有工程重擔，「領域專業知識」成為最高維度的生產力",
        category="11 · CONCLUSION & FUTURE"
    )
    add_card(
        slide12, Inches(0.8), Inches(1.85), Inches(5.6), Inches(5.1),
        "🚀 本專案帶給我們的核心啟示",
        [
            "1. **領域專家就是最佳架構師**：醫師最懂臨床流程，在 Vibe Coding 下，醫師能直接主導軟體進化，省去與外包工程師數月的溝通代溝。",
            "2. **極致敏捷的創新閉環**：一個醫療點子（如自訂日期警示）從提出、法規檢驗、編譯修復到 Android/Web 發布，僅耗時數分鐘。",
            "3. **架構思維 > 語法記憶**：未來的核心競爭力是「清晰的需求定義能力」與「臨床批判性審查」。"
        ],
        badge="Key Takeaways",
        header_color=C_BLUE
    )
    add_card(
        slide12, Inches(6.8), Inches(1.85), Inches(5.7), Inches(5.1),
        "🌟 未來展望：人人皆是數位創造者",
        [
            "• 醫護人員可以為自己的診所打造專屬臨床工具系統。",
            "• 教師可以為學生打造即時互動測驗與學習進度追蹤器。",
            "• 企業主可以在一天內驗證並交付跨平台商業系統。",
            "",
            "✨ 『在 Antigravity 與 Vibe Coding 的時代，你的想像力與專業知識，就是最強大的程式語言。』"
        ],
        badge="Empowerment",
        border_color=RGBColor(254, 215, 170),
        header_color=C_AMBER
    )

    output_path = r"E:\Vaccine\Vibe_Coding_Presentation.pptx"
    prs.save(output_path)
    print(f"Authentic Desktop UI 12-slide presentation successfully generated: {output_path}")

if __name__ == "__main__":
    build_final_authentic_presentation()
