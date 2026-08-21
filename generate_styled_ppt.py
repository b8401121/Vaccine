import os
import math
import numpy as np
from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 1. 產生超高清 16:9 醫療科技風格質感背景底圖 (1920x1080)
def generate_medical_tech_backgrounds():
    width, height = 1920, 1080
    bg_dir = r"E:\Vaccine\ppt_assets"
    os.makedirs(bg_dir, exist_ok=True)
    
    # 封面背景 (Cover Background: 深度漸層 + 數位 DNA 粒子波浪 + 科技網格)
    cover_img = Image.new("RGBA", (width, height), (11, 19, 43, 255))
    draw = ImageDraw.Draw(cover_img)
    
    # 徑向/線性深色漸層
    for y in range(height):
        r = int(11 + (22 - 11) * (y / height))
        g = int(19 + (38 - 19) * (y / height))
        b = int(43 + (70 - 43) * (y / height))
        draw.line([(0, y), (width, y)], fill=(r, g, b, 255))
        
    # 繪製抽象流動曲線 (DNA 波浪 / 生長曲線意象)
    overlay = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    d_overlay = ImageDraw.Draw(overlay)
    
    for x in range(0, width, 4):
        # 曲線 1 (Cyan)
        y1 = int(height * 0.72 + math.sin(x * 0.0035) * 110 + math.cos(x * 0.007) * 45)
        d_overlay.ellipse([x-2, y1-2, x+2, y1+2], fill=(56, 189, 248, 60))
        
        # 曲線 2 (Teal)
        y2 = int(height * 0.78 + math.cos(x * 0.004) * 90 + math.sin(x * 0.008) * 35)
        d_overlay.ellipse([x-2, y2-2, x+2, y2+2], fill=(45, 212, 191, 55))
        
        # 曲線 3 (Deep Blue)
        y3 = int(height * 0.65 + math.sin(x * 0.005 + 1) * 80)
        d_overlay.ellipse([x-1, y3-1, x+1, y3+1], fill=(129, 140, 248, 40))

    # 右上角抽象發光微粒 / 疫苗抗體節點網 (Nodes)
    np.random.seed(42)
    node_points = []
    for _ in range(35):
        nx = int(np.random.uniform(width * 0.65, width * 0.98))
        ny = int(np.random.uniform(height * 0.05, height * 0.45))
        node_points.append((nx, ny))
        d_overlay.ellipse([nx-4, ny-4, nx+4, ny+4], fill=(56, 189, 248, 140))
        d_overlay.ellipse([nx-8, ny-8, nx+8, ny+8], fill=(56, 189, 248, 40))

    # 連接相鄰節點
    for i, p1 in enumerate(node_points):
        for j, p2 in enumerate(node_points):
            if i < j:
                dist = math.hypot(p1[0] - p2[0], p1[1] - p2[1])
                if dist < 160:
                    alpha = int((1 - dist / 160) * 50)
                    d_overlay.line([p1, p2], fill=(45, 212, 191, alpha), width=1)

    # 模糊光暈效果疊加
    glow = overlay.filter(ImageFilter.GaussianBlur(radius=3))
    cover_final = Image.alpha_composite(cover_img, glow)
    cover_final = Image.alpha_composite(cover_final, overlay)
    
    cover_path = os.path.join(bg_dir, "cover_bg.png")
    cover_final.convert("RGB").save(cover_path, quality=95)
    
    # 內頁背景 (Slide Background: 雅緻科技深色微紋理，不干擾文字閱讀)
    slide_img = Image.new("RGBA", (width, height), (15, 23, 42, 255))
    d_slide = ImageDraw.Draw(slide_img)
    for y in range(height):
        r = int(13 + (20 - 13) * (y / height))
        g = int(20 + (30 - 20) * (y / height))
        b = int(38 + (55 - 38) * (y / height))
        d_slide.line([(0, y), (width, y)], fill=(r, g, b, 255))

    slide_overlay = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    d_so = ImageDraw.Draw(slide_overlay)

    # 底部微光科技波紋
    for x in range(0, width, 6):
        sy = int(height * 0.94 + math.sin(x * 0.005) * 35)
        d_so.ellipse([x-2, sy-2, x+2, sy+2], fill=(56, 189, 248, 35))
        
    # 右上角極輕量網格點陣
    for gx in range(int(width * 0.78), width, 32):
        for gy in range(0, int(height * 0.25), 32):
            d_so.ellipse([gx-1, gy-1, gx+1, gy+1], fill=(148, 163, 184, 25))

    slide_glow = slide_overlay.filter(ImageFilter.GaussianBlur(radius=2))
    slide_final = Image.alpha_composite(slide_img, slide_glow)
    slide_final = Image.alpha_composite(slide_final, slide_overlay)
    
    content_bg_path = os.path.join(bg_dir, "content_bg.png")
    slide_final.convert("RGB").save(content_bg_path, quality=95)
    
    return cover_path, content_bg_path


# 2. 建立精美簡報
def build_styled_presentation():
    cover_bg, content_bg = generate_medical_tech_backgrounds()
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # 配色常數
    COLOR_CARD = RGBColor(30, 41, 59)          # #1E293B 卡片底色 (半透感質感)
    COLOR_CARD_BORDER = RGBColor(51, 65, 85)   # #334155 卡片微亮邊框
    COLOR_PRIMARY = RGBColor(56, 189, 248)     # #38BDF8 天空藍
    COLOR_TEAL = RGBColor(45, 212, 191)        # #2DD4BF 醫療薄荷綠
    COLOR_AMBER = RGBColor(251, 191, 36)       # #FBBF24 警示橘金
    COLOR_TEXT_MAIN = RGBColor(248, 250, 252)  # #F8FAFC 主文字白
    COLOR_TEXT_MUTED = RGBColor(203, 213, 225) # #CBD5E1 高對比清晰說明灰

    def add_base_slide(title_text, subtitle_text=None, category_badge="VIBE CODING ARCHITECTURE"):
        slide = prs.slides.add_slide(blank_layout)
        
        # 加入高清質感背景圖
        slide.shapes.add_picture(content_bg, 0, 0, prs.slide_width, prs.slide_height)

        # 頂部分類小膠囊標籤
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(2.8), Inches(0.3))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(15, 45, 75)
        badge.line.color.rgb = COLOR_PRIMARY
        badge.line.width = Pt(1)
        
        tb_b = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(2.8), Inches(0.3))
        tf_b = tb_b.text_frame
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
        p_b = tf_b.paragraphs[0]
        p_b.text = f"◈ {category_badge}"
        p_b.font.name = "Microsoft JhengHei"
        p_b.font.size = Pt(9.5)
        p_b.font.bold = True
        p_b.font.color.rgb = COLOR_PRIMARY

        # 標題文字方塊
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.78), Inches(11.7), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Microsoft JhengHei"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = "Microsoft JhengHei"
            p2.font.size = Pt(12)
            p2.font.color.rgb = COLOR_TEXT_MUTED
            p2.space_before = Pt(3)

        return slide

    def add_card(slide, left, top, width, height, title, content_list, badge_text=None, border_color=COLOR_CARD_BORDER, header_color=COLOR_PRIMARY):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)

        tb = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.2), width - Inches(0.44), height - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Microsoft JhengHei"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = header_color

        if badge_text:
            p.text = f"{title}  [{badge_text}]"

        for item in content_list:
            p_item = tf.add_paragraph()
            p_item.text = item
            p_item.font.name = "Microsoft JhengHei"
            p_item.font.size = Pt(11.5)
            p_item.font.color.rgb = COLOR_TEXT_MUTED
            p_item.space_before = Pt(5)
            if item.startswith("•") or item.startswith("👉") or item.startswith("✓") or item.startswith("🚨") or item.startswith("【"):
                p_item.font.color.rgb = COLOR_TEXT_MAIN

    # -------------------------------------------------------------
    # 投影片 1：封面 (Cover Slide)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    slide1.shapes.add_picture(cover_bg, 0, 0, prs.slide_width, prs.slide_height)

    # 頂部專案小標籤
    tag_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.3), Inches(3.8), Inches(0.38))
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = RGBColor(12, 74, 110)
    tag_box.line.color.rgb = COLOR_PRIMARY
    tag_box.line.width = Pt(1)
    
    t_tag = tag_box.text_frame
    p_t = t_tag.paragraphs[0]
    p_t.text = "✦ AI-NATIVE SOFTWARE ENGINEERING"
    p_t.font.name = "Microsoft JhengHei"
    p_t.font.size = Pt(11)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_PRIMARY

    tb1 = slide1.shapes.add_textbox(Inches(1.2), Inches(1.9), Inches(11.0), Inches(4.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "何謂 Vibe Coding？"
    p1.font.name = "Microsoft JhengHei"
    p1.font.size = Pt(46)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_MAIN

    p2 = tf1.add_paragraph()
    p2.text = "從「台灣預防接種指南助手」看 AI 時代人機共創的全新軟體開發典範"
    p2.font.name = "Microsoft JhengHei"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_PRIMARY
    p2.space_before = Pt(10)

    p3 = tf1.add_paragraph()
    p3.text = "實務專案案例剖析 ‧ 核心概念演進 ‧ 人機角色蛻變 ‧ 醫療輔助系統四端跨平台落地"
    p3.font.name = "Microsoft JhengHei"
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_TEAL
    p3.space_before = Pt(22)

    p4 = tf1.add_paragraph()
    p4.text = "案例系統：Rust 醫療算法 ‧ WebAssembly 純前端運算 ‧ Android 原生 APK ‧ iOS 原生日曆連動"
    p4.font.name = "Microsoft JhengHei"
    p4.font.size = Pt(12)
    p4.font.color.rgb = COLOR_TEXT_MUTED
    p4.space_before = Pt(8)

    # -------------------------------------------------------------
    # 投影片 2：什麼是 Vibe Coding？
    # -------------------------------------------------------------
    slide2 = add_base_slide(
        "1. 什麼是 Vibe Coding？核心定義與概念演進",
        "由前 Tesla AI 總監、OpenAI 共同創辦人 Andrej Karpathy 於 2025 年初提出的顛覆性開發思維",
        category_badge="PARADIGM SHIFT"
    )
    add_card(
        slide2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.1),
        "💡 Vibe Coding 的核心本質",
        [
            "• 開發者不再逐字手寫語法，而是以「直覺、意圖與自然語言對話」引導 AI。",
            "• 開發者的角色蛻變為【總架構師】與【產品質檢官 (Director & Reviewer)】。",
            "• 重點放在「想要什麼體驗 (Vibe) 與核心業務邏輯」，讓 AI Agent 負責底層實作、修復報錯、編譯與跨平台部署。",
            "👉 『You just see stuff, say stuff, run stuff, and copy-paste stuff. It works.』"
        ],
        badge_text="全新心智模型",
        header_color=COLOR_PRIMARY
    )
    add_card(
        slide2, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.1),
        "⚡ 傳統開發 vs. Vibe Coding 對比",
        [
            "【傳統工程模式 (Traditional Coding)】",
            "• 手動查閱 API 文件、逐行記憶並撰寫 boilerplate 樣板代碼",
            "• 遭遇跨平台 (Web / Android / NDK) 編譯地獄時耗費數天排錯",
            "",
            "【Vibe Coding 人機協作模式】",
            "• 領域專家（如醫師/PM）直接以專業業務邏輯下達指令",
            "• AI 即刻進行法規文獻查證、Rust 算法撰寫、Wasm 封裝、APK 簽署",
            "• 敏捷迭代：一個複雜臨床需求在數分鐘內完成跨平台發布"
        ],
        border_color=COLOR_TEAL,
        header_color=COLOR_TEAL
    )

    # -------------------------------------------------------------
    # 投影片 3：經典案例剖析 - 台灣預防接種指南助手
    # -------------------------------------------------------------
    slide3 = add_base_slide(
        "2. 實戰案例剖析：台灣預防接種指南助手",
        "以真實複雜的醫療系統為例——如何透過自然語言對話完成四端跨平台同步開發",
        category_badge="REAL-WORLD CASE STUDY"
    )
    add_card(
        slide3, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5.1),
        "🏥 臨床專業意圖 (Doctor's Intent)",
        [
            "• 「診所需要算公自費疫苗時程」",
            "• 「7歲以上兒童有生長曲線嗎？請引用台灣最新官方年份」",
            "• 「要能支援 iOS 相機直掃行事曆」",
            "• 「家長改日期若提早，需跳出醫療法規理由警示」",
            "",
            "👉 使用者專注於【醫療核心邏輯】，完全不需觸碰底層語法。"
        ],
        badge_text="Prompt as Logic",
        header_color=COLOR_PRIMARY
    )
    add_card(
        slide3, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5.1),
        "🤖 AI Agent 自主架構與執行",
        [
            "• 聯網查核 2024 衛福部國健署標準",
            "• 撰寫嚴謹的 Rust 演算法 (vaccine-core)",
            "• 編譯 WebAssembly 驅動網頁純前端運算",
            "• 自動配置 Android NDK / Clang 工具鏈打包簽署 APK",
            "• 生成 iCalendar (.ics) / VEVENT 數據流"
        ],
        badge_text="Autonomous Agent",
        border_color=COLOR_PRIMARY,
        header_color=COLOR_TEAL
    )
    add_card(
        slide3, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.1),
        "🚀 四端跨平台成果 (Deliverables)",
        [
            "1. 🌐 GitHub Pages 靜態無伺服器網站",
            "2. 📱 Android 原生 arm64 簽署 APK",
            "3. 💻 Windows 單一免安裝獨立綠色版",
            "4. 🍏 iOS 相機直掃 iCal 日曆提醒",
            "",
            "✨ 從單一需求到四端落地，全由對話驅動！"
        ],
        badge_text="Multi-Platform",
        border_color=COLOR_TEAL,
        header_color=COLOR_AMBER
    )

    # -------------------------------------------------------------
    # 投影片 4：Vibe Coding 運作流程
    # -------------------------------------------------------------
    slide4 = add_base_slide(
        "3. Vibe Coding 的四步驟高頻迭代循環",
        "從自然語言需求到跨平台生產環境交付的完整閉環 (Feedback Loop)",
        category_badge="ITERATIVE WORKFLOW"
    )
    steps = [
        ("Step 1: 意圖表達", "使用者提出自然語言需求\n例如：「增加iOS行事曆QR Code與自訂日期」", COLOR_PRIMARY),
        ("Step 2: 查核與架構", "AI 搜尋法規、規劃資料模型\n設計前後端跨語言序列化格式", COLOR_TEAL),
        ("Step 3: 實現與編譯", "AI 撰寫程式碼、呼叫編譯器\n遇到報錯自主排查修復 (Self-Healing)", COLOR_AMBER),
        ("Step 4: 驗證與發布", "產出 Android APK、推送到 GitHub\n使用者即刻在手機/網頁體驗", COLOR_PRIMARY),
    ]
    for i, (title, desc, color) in enumerate(steps):
        left = Inches(0.8 + i * 2.95)
        add_card(
            slide4, left, Inches(1.9), Inches(2.8), Inches(4.9),
            title,
            [desc, "", "👉 立即進入下一輪反饋循環 (Feedback Loop)"],
            border_color=color,
            header_color=color
        )

    # -------------------------------------------------------------
    # 投影片 5：關鍵支柱與成功要素
    # -------------------------------------------------------------
    slide5 = add_base_slide(
        "4. Vibe Coding 成功落地的三大支柱",
        "為什麼本專案能做到又快又穩？底層的技術架構與工程支撐",
        category_badge="CORE PILLARS"
    )
    add_card(
        slide5, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5.1),
        "1. Skill 與上下文沉澱",
        [
            "• 建立專屬 SKILL.md 開發規範",
            "• 將 NDK 環境變數、Wasm 序列化規範、Git 部署指令結構化",
            "• 讓 AI 每次對話都能維持統一水準，不遺忘專案背景知識"
        ],
        badge_text="Persistent Memory",
        header_color=COLOR_PRIMARY
    )
    add_card(
        slide5, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5.1),
        "2. 工具調用能力 (Tool Calling)",
        [
            "• 具備終端機 (PowerShell)、網路搜尋、檔案精確替換工具",
            "• 遇到 Android 編譯錯誤時，AI 能自主讀取報錯並修正 Rust / Gradle 配置",
            "• 真正的【自主除錯與修復 (Self-Healing)】"
        ],
        badge_text="Autonomous Agent",
        border_color=COLOR_PRIMARY,
        header_color=COLOR_TEAL
    )
    add_card(
        slide5, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.1),
        "3. 領域專業引導 (Domain Knowledge)",
        [
            "• 人類的不可替代性在於【臨床專業判斷】",
            "• 審核疫苗間隔邏輯是否合規、評估使用者互動體驗",
            "• 領域專家 + 強大 AI = 專業級系統迅速誕生"
        ],
        badge_text="Human-in-the-Loop",
        border_color=COLOR_TEAL,
        header_color=COLOR_AMBER
    )

    # -------------------------------------------------------------
    # 投影片 6：總結與未來展望
    # -------------------------------------------------------------
    slide6 = add_base_slide(
        "5. 總結：軟體開發的文藝復興時代",
        "每個人都能成為自己領域的數位架構師 (Domain Architects)",
        category_badge="FUTURE VISION"
    )
    add_card(
        slide6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.1),
        "🚀 核心啟示 (Key Takeaways)",
        [
            "1. 程式碼不再是壁壘：思想與領域知識 (Domain Knowledge) 才是核心價值。",
            "2. 敏捷的極致體現：從醫療想法到多端上線僅需數小時，大幅降低創新試錯成本。",
            "3. 架構思維 > 語法記憶：掌握系統設計與需求拆解能力，比死記 API 更有競爭力。"
        ],
        badge_text="Summary",
        header_color=COLOR_PRIMARY
    )
    add_card(
        slide6, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.1),
        "🌟 未來展望 (The Future of Creation)",
        [
            "• 醫生可以為診所打造最合身的醫療輔助系統",
            "• 老師可以為學生建構客製化的互動教學軟體",
            "• 創業家可以在一天內驗證並發布 MVP 產品",
            "",
            "✨ 『在 AI 時代，你的想像力與專業知識，就是最強大的程式語言。』"
        ],
        border_color=COLOR_AMBER,
        header_color=COLOR_AMBER
    )

    output_path = r"E:\Vaccine\Vibe_Coding_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation generated successfully: {output_path}")

if __name__ == "__main__":
    build_styled_presentation()
